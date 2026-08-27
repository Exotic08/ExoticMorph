"""Tests for the Narrative Arc (3-part) layout system.

Covers the v7 upgrade:
- `LayoutType` gains the designer layouts COVER_HERO / BIG_STAT_CALLOUT /
  ASYMMETRIC_GRID / CONCLUSION_SUMMARY (CARDS_ROW & TIMELINE_STEPS kept) and the
  legacy values still parse (backwards compatible).
- `layout_engine._resolve_slide_layouts()` ENFORCES the narrative arc:
  slide 1 = COVER_HERO, slide N = CONCLUSION_SUMMARY, body slides alternate
  through BODY_LAYOUTS and never repeat the previous slide's layout.
- Geometry contract for each new layout (font sizes, 60/40 split, centred
  three thin summary rows, nothing outside the 16:9 canvas, no overlapping cards).
- The Pydantic schema still round-trips to JSON (frontend preview contract) and
  the few-shot examples embedded in the system prompt are schema-valid.
"""

from __future__ import annotations

import io
import json
import os
import sys
import unittest

BACKEND_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if BACKEND_ROOT not in sys.path:
    sys.path.insert(0, BACKEND_ROOT)

from app.schemas.slide_schema import (  # noqa: E402
    BODY_LAYOUTS,
    CLOSING_LAYOUTS,
    GenerateRequest,
    LEGACY_LAYOUT_UPGRADE,
    LLMOutput,
    MAX_ELEMENTS_PER_SLIDE,
    NARRATIVE_LAYOUTS,
    OPENING_LAYOUTS,
    PresentationResponse,
    SlideData,
    SlideElement,
    LayoutType,
    _coerce_layout_type,
)
from app.services import layout_engine as le  # noqa: E402
from app.services.layout_engine import (  # noqa: E402
    _calculate_cards_geometry_inches,
    _conclusion_frame_geometry,
    _conclusion_rows_geometry,
    _cover_geometry,
    _cover_title_font_size,
    _resolve_slide_layouts,
    _stat_number_font_size,
    compute_layout,
)
from app.services.llm_service import (  # noqa: E402
    _FEW_SHOT_EXAMPLE_1,
    _FEW_SHOT_EXAMPLE_2,
    SYSTEM_PROMPT,
    _build_user_prompt,
)

# 16:9 canvas in inches (mirrors layout_engine constants)
SLIDE_W = le.SLIDE_W_IN
SLIDE_H = le.SLIDE_H_IN


# ==============================================================================
# FIXTURES
# ==============================================================================

def _card(i: int, title: str, desc: str, color: str = "#8B5CF6") -> dict:
    return {
        "morph_id": f"card_{i}",
        "title": title,
        "description": desc,
        "color_theme": color,
        "order": i,
    }


DESC = (
    "Khí quyển Kim Tinh dày 92 bar với 96.5% CO2 tạo hiệu ứng nhà kính mất kiểm "
    "soát, giữ nhiệt độ bề mặt ở 465 độ C cả ngày lẫn đêm."
)


def _deck_payload(num_slides: int = 5, layout: str = "CARDS_ROW") -> dict:
    """A schema-valid deck where EVERY slide uses the same (body) layout."""
    slides = []
    for i in range(1, num_slides + 1):
        slides.append({
            "slide_number": i,
            "section": f"PHẦN {i}",
            "slide_title": f"Luận điểm số {i} về khí quyển Kim Tinh và CO2",
            "layout_type": layout,
            "cards": [
                _card(0, f"Luận cứ {i} về áp suất 92 bar", DESC, "#8B5CF6"),
                _card(1, f"Luận cứ {i} về nhiệt độ 465 độ C", DESC, "#10B981"),
                _card(2, f"Luận cứ {i} về mây axit sunfuric", DESC, "#06B6D4"),
            ],
        })
    return {"topic": "Hệ Mặt Trời", "slides": slides}


def _arc_payload() -> dict:
    """A deck written the way the system prompt asks for (correct arc + card roles)."""
    return {
        "topic": "Hệ Mặt Trời",
        "slides": [
            {
                "slide_number": 1,
                "section": "CÂU HỎI MỞ ĐẦU",
                "slide_title": "Vì sao Kim Tinh nóng 465 °C dù xa Mặt Trời hơn Thuỷ Tinh?",
                "layout_type": "COVER_HERO",
                "cards": [
                    _card(0, "Một nghịch lý nhiệt động giữa hai hành tinh láng giềng", DESC),
                    _card(1, "KHÍ QUYỂN", "Lớp vỏ khí 92 bar chi phối nhiệt động học.", "#10B981"),
                    _card(2, "QUỸ ĐẠO", "Chu kỳ 225 ngày, gần tròn nhất hệ.", "#06B6D4"),
                ],
            },
            {
                "slide_number": 2,
                "section": "SỐ LIỆU",
                "slide_title": "Áp suất bề mặt Kim Tinh gấp 92 lần Trái Đất",
                "layout_type": "BIG_STAT_CALLOUT",
                "cards": [
                    _card(0, "92 bar", "Tương đương độ sâu 900 m dưới đại dương Trái Đất."),
                    _card(1, "CO2 chiếm 96.5%", "Khí nhà kính chính, mây H2SO4 phản xạ 75% ánh sáng.", "#10B981"),
                ],
            },
            {
                "slide_number": 3,
                "section": "CƠ CHẾ",
                "slide_title": "Ba tầng khí quyển vận hành như một nồi áp suất",
                "layout_type": "ASYMMETRIC_GRID",
                "cards": [
                    _card(0, "Tầng đối lưu dày 65 km giữ trọn nhiệt lượng", DESC),
                    _card(1, "Mây axit sunfuric", "Phản xạ ánh sáng nhưng bẫy bức xạ hồng ngoại.", "#10B981"),
                    _card(2, "Hiệu ứng nhà kính", "CO2 và SO2 tạo vòng phản hồi không đảo ngược.", "#06B6D4"),
                ],
            },
            {
                "slide_number": 4,
                "section": "TIẾN TRÌNH",
                "slide_title": "Bốn giai đoạn hình thành nghịch lý nhiệt",
                "layout_type": "TIMELINE_STEPS",
                "cards": [
                    _card(0, "Giai đoạn 1: đại dương bốc hơi", "Hơi nước bổ sung vào khí quyển nguyên thuỷ."),
                    _card(1, "Giai đoạn 2: CO2 tích tụ", "Núi lửa giải phóng CO2 không bị phong hoá hấp thụ.", "#10B981"),
                    _card(2, "Giai đoạn 3: mất nước", "UV phân ly H2O, hydro thoát ra không gian.", "#06B6D4"),
                    _card(3, "Giai đoạn 4: cân bằng 465 °C", "Khí quyển đạt trạng thái siêu tới hạn ổn định.", "#EC4899"),
                ],
            },
            {
                "slide_number": 5,
                "section": "KẾT LUẬN",
                "slide_title": "Kim Tinh là lời cảnh báo cho Trái Đất",
                "layout_type": "CONCLUSION_SUMMARY",
                "cards": [
                    _card(0, "Khoảng cách không quyết định nhiệt độ", "Thành phần khí quyển mới là biến số chi phối."),
                    _card(1, "Hiệu ứng nhà kính có điểm không quay lại", "Vượt ngưỡng bốc hơi đại dương là irreversible.", "#10B981"),
                    _card(2, "Trái Đất đang ở rìa vùng an toàn", "Mô hình khí hậu cho thấy +4 °C kích hoạt vòng phản hồi.", "#06B6D4"),
                    _card(3, "Theo dõi nồng độ CO2 toàn cầu ngay hôm nay", "Giám sát ppm CO2 tại Mauna Loa và cắt giảm phát thải.", "#EC4899"),
                ],
            },
        ],
    }


def _compute(payload: dict, num_slides: int | None = None, style: str = "Futuristic"):
    raw = LLMOutput.model_validate(payload)
    req = GenerateRequest(
        prompt="hệ mặt trời",
        num_slides=num_slides or len(payload["slides"]),
        style=style,
    )
    return compute_layout(raw, req)


def _elems(slide, *types):
    return [e for e in slide.elements if e.type in types]


def _overlap(a: SlideElement, b: SlideElement) -> bool:
    return not (
        a.x + a.width <= b.x + 1e-6
        or b.x + b.width <= a.x + 1e-6
        or a.y + a.height <= b.y + 1e-6
        or b.y + b.height <= a.y + 1e-6
    )


# ==============================================================================
# 1. SCHEMA — LayoutType enum, aliases, narrative groups
# ==============================================================================

class LayoutTypeEnumTests(unittest.TestCase):
    def test_new_designer_layouts_exist(self) -> None:
        for name in (
            "COVER_HERO", "BIG_STAT_CALLOUT", "ASYMMETRIC_GRID",
            "CARDS_ROW", "TIMELINE_STEPS", "CONCLUSION_SUMMARY",
        ):
            self.assertIn(name, LayoutType.__members__)
            self.assertEqual(LayoutType[name].value, name)

    def test_legacy_layouts_still_supported(self) -> None:
        for name in ("SPLIT_HERO", "STAT_GRID", "GRID_2X2"):
            self.assertIn(name, LayoutType.__members__)

    def test_narrative_groups_are_disjoint_and_complete(self) -> None:
        self.assertEqual(OPENING_LAYOUTS, (LayoutType.COVER_HERO,))
        self.assertEqual(CLOSING_LAYOUTS, (LayoutType.CONCLUSION_SUMMARY,))
        self.assertEqual(
            BODY_LAYOUTS,
            (
                LayoutType.BIG_STAT_CALLOUT,
                LayoutType.ASYMMETRIC_GRID,
                LayoutType.TIMELINE_STEPS,
                LayoutType.CARDS_ROW,
            ),
        )
        groups = list(OPENING_LAYOUTS) + list(BODY_LAYOUTS) + list(CLOSING_LAYOUTS)
        self.assertEqual(len(groups), len(set(groups)))
        self.assertEqual(set(NARRATIVE_LAYOUTS), set(groups))
        for legacy in LEGACY_LAYOUT_UPGRADE:
            self.assertNotIn(legacy, NARRATIVE_LAYOUTS)
            self.assertIn(LEGACY_LAYOUT_UPGRADE[legacy], BODY_LAYOUTS)

    def test_alias_coercion(self) -> None:
        cases = {
            "cover_hero": LayoutType.COVER_HERO,
            "Cover": LayoutType.COVER_HERO,
            "title-hero": LayoutType.COVER_HERO,
            "opening": LayoutType.COVER_HERO,
            "big stat": LayoutType.BIG_STAT_CALLOUT,
            "Stat_Callout": LayoutType.BIG_STAT_CALLOUT,
            "asymmetric-grid": LayoutType.ASYMMETRIC_GRID,
            "ASYM_GRID": LayoutType.ASYMMETRIC_GRID,
            "takeaways": LayoutType.CONCLUSION_SUMMARY,
            "Conclusion": LayoutType.CONCLUSION_SUMMARY,
            "summary_card": LayoutType.CONCLUSION_SUMMARY,
            "timeline": LayoutType.TIMELINE_STEPS,
        }
        for raw, expected in cases.items():
            self.assertEqual(_coerce_layout_type(raw), expected, f"alias {raw!r}")

    def test_unknown_layout_falls_back_to_cards_row(self) -> None:
        self.assertEqual(_coerce_layout_type("KHONG_TON_TAI"), LayoutType.CARDS_ROW)
        self.assertEqual(_coerce_layout_type(None), LayoutType.CARDS_ROW)

    def test_slide_data_accepts_new_and_legacy_layouts(self) -> None:
        for layout in ("COVER_HERO", "BIG_STAT_CALLOUT", "CONCLUSION_SUMMARY", "GRID_2X2"):
            slide = SlideData.model_validate({
                "slide_number": 1,
                "section": "TEST",
                "slide_title": "Tiêu đề kiểm thử",
                "layout_type": layout,
                "cards": [_card(0, "Thẻ kiểm thử số một", DESC)],
            })
            self.assertEqual(slide.layout_type.value, layout)

    def test_llm_schema_still_rejects_geometry(self) -> None:
        payload = _arc_payload()
        payload["slides"][0]["cards"][0]["x"] = 12.5
        with self.assertRaises(Exception):
            LLMOutput.model_validate(payload)


# ==============================================================================
# 2. NARRATIVE ARC ENFORCEMENT
# ==============================================================================

class NarrativeArcResolverTests(unittest.TestCase):
    def _resolve(self, payload: dict) -> list[LayoutType]:
        raw = LLMOutput.model_validate(payload)
        return _resolve_slide_layouts(raw.slides)

    def test_single_slide_is_cover_hero(self) -> None:
        layouts = self._resolve(_deck_payload(1))
        self.assertEqual(layouts, [LayoutType.COVER_HERO])

    def test_two_slides_open_and_close(self) -> None:
        layouts = self._resolve(_deck_payload(2))
        self.assertEqual(layouts, [LayoutType.COVER_HERO, LayoutType.CONCLUSION_SUMMARY])

    def test_five_slides_follow_arc_even_when_llm_repeats_layout(self) -> None:
        layouts = self._resolve(_deck_payload(5, layout="CARDS_ROW"))
        self.assertEqual(layouts[0], LayoutType.COVER_HERO)
        self.assertEqual(layouts[-1], LayoutType.CONCLUSION_SUMMARY)
        for layout in layouts[1:-1]:
            self.assertIn(layout, BODY_LAYOUTS)
        for prev, cur in zip(layouts, layouts[1:]):
            self.assertNotEqual(prev, cur, f"2 slide liên tiếp trùng layout {prev}")

    def test_ten_slides_never_repeat_consecutive_layout(self) -> None:
        layouts = self._resolve(_deck_payload(10, layout="BIG_STAT_CALLOUT"))
        self.assertEqual(layouts[0], LayoutType.COVER_HERO)
        self.assertEqual(layouts[-1], LayoutType.CONCLUSION_SUMMARY)
        for prev, cur in zip(layouts, layouts[1:]):
            self.assertNotEqual(prev, cur)
        for layout in layouts[1:-1]:
            self.assertIn(layout, BODY_LAYOUTS)
        # 8 slide thân bài nhưng chỉ có 4 layout → phải luân phiên, không được trùng kề
        self.assertGreaterEqual(len({l.value for l in layouts[1:-1]}), 2)

    def test_cover_and_conclusion_never_appear_in_body(self) -> None:
        layouts = self._resolve(_deck_payload(7, layout="COVER_HERO"))
        self.assertNotIn(LayoutType.COVER_HERO, layouts[1:])
        self.assertNotIn(LayoutType.CONCLUSION_SUMMARY, layouts[1:-1])

    def test_legacy_layouts_are_upgraded(self) -> None:
        for legacy, upgraded in LEGACY_LAYOUT_UPGRADE.items():
            layouts = self._resolve(_deck_payload(5, layout=legacy.value))
            for layout in layouts[1:-1]:
                self.assertIn(layout, BODY_LAYOUTS)
                self.assertNotIn(layout, LEGACY_LAYOUT_UPGRADE)
            self.assertIn(upgraded, BODY_LAYOUTS)

    def test_correct_llm_arc_is_preserved(self) -> None:
        """Khi LLM làm đúng, engine KHÔNG đổi layout của nó."""
        layouts = self._resolve(_arc_payload())
        self.assertEqual(
            [l.value for l in layouts],
            ["COVER_HERO", "BIG_STAT_CALLOUT", "ASYMMETRIC_GRID",
             "TIMELINE_STEPS", "CONCLUSION_SUMMARY"],
        )


# ==============================================================================
# 3. GEOMETRY — COVER_HERO
# ==============================================================================

class CoverHeroGeometryTests(unittest.TestCase):
    def test_text_line_estimator_counts_long_unbroken_tokens(self) -> None:
        # Một chuỗi 190 ký tự không có khoảng trắng PHẢI được tính nhiều dòng.
        self.assertGreater(le._estimate_text_lines("k" * 190, 44, 9.8), 3)
        self.assertEqual(le._estimate_text_lines("", 44, 9.8), 1)
        self.assertEqual(le._estimate_text_lines("Ngắn", 44, 9.8), 1)
        self.assertGreater(
            le._estimate_text_lines("từ " * 60, 44, 9.8),
            le._estimate_text_lines("từ " * 10, 44, 9.8),
        )

    def test_title_font_size_within_designer_range(self) -> None:
        for title in ("Ngắn", "x" * 40, "y" * 70, "z" * 150):
            size = _cover_title_font_size(title)
            self.assertGreaterEqual(size, 44)
            self.assertLessEqual(size, 52)
        self.assertEqual(_cover_title_font_size("Ngắn gọn"), 52)
        self.assertEqual(_cover_title_font_size("z" * 150), 44)

    def test_cover_stack_is_ordered_and_inside_canvas(self) -> None:
        geo = _cover_geometry(
            "Vì sao Kim Tinh nóng 465 °C dù xa Mặt Trời hơn Thuỷ Tinh?",
            subtitle="Một nghịch lý nhiệt động giữa hai hành tinh láng giềng",
            note=DESC,
            badge_labels=["KHÍ QUYỂN", "QUỸ ĐẠO", "NHIỆT ĐỘ"],
        )
        tx, ty, tw, th = geo["title"]
        rx, ry, rw, rh = geo["rule"]
        sx, sy, sw, sh = geo["subtitle"]

        self.assertEqual(tx, le.COVER_MX_IN)
        self.assertEqual(rx, le.COVER_MX_IN)
        self.assertGreater(ry, ty + th, "thanh accent phải nằm dưới tiêu đề")
        self.assertGreater(sy, ry + rh, "subtitle phải nằm dưới thanh accent")
        self.assertLessEqual(ty + th, SLIDE_H)
        self.assertLessEqual(sy + sh, SLIDE_H)
        self.assertLessEqual(tx + tw, SLIDE_W)

        self.assertEqual(len(geo["badges"]), 3)
        prev_right = None
        for bx, by, bw, bh in geo["badges"]:
            self.assertGreater(by, sy + sh, "badge phải nằm dưới subtitle")
            self.assertLessEqual(by + bh, le.FOOTER_TOP_IN, "badge chạm footer")
            self.assertLessEqual(bx + bw, SLIDE_W - le.COVER_MX_IN + 1e-6)
            if prev_right is not None:
                self.assertGreater(bx, prev_right, "các badge chồng lấn nhau")
            prev_right = bx + bw

    def test_long_title_never_pushes_badges_into_footer(self) -> None:
        geo = _cover_geometry(
            "k" * 190,
            subtitle="s" * 200,
            badge_labels=["NHÃN RẤT DÀI SỐ MỘT", "NHÃN RẤT DÀI SỐ HAI",
                          "NHÃN RẤT DÀI SỐ BA", "NHÃN RẤT DÀI SỐ BỐN"],
        )
        for _, by, _, bh in geo["badges"]:
            self.assertLessEqual(by + bh, le.FOOTER_TOP_IN)

    def test_wrapping_badges_are_dropped_rather_than_hitting_the_footer(self) -> None:
        """4 badge dài (wrap 2 hàng) + tiêu đề/subtitle cực dài → phải bỏ bớt badge."""
        geo = _cover_geometry(
            "k" * 190,
            subtitle="s" * 200,
            badge_labels=["A" * 32, "B" * 32, "C" * 32, "D" * 32],
        )
        self.assertLess(len(geo["badges"]), 4)
        for _, by, _, bh in geo["badges"]:
            self.assertLessEqual(by + bh, le.FOOTER_TOP_IN)
        if geo["subtitle"]:
            sy, sh = geo["subtitle"][1], geo["subtitle"][3]
            for _, by, _, _ in geo["badges"]:
                self.assertGreater(by, sy + sh)

    def test_single_card_uses_note_instead_of_badges(self) -> None:
        geo = _cover_geometry("Tiêu đề mở bài", subtitle="Câu hook", note=DESC, badge_labels=[])
        self.assertEqual(geo["badges"], [])
        self.assertIsNotNone(geo["note"])

    def test_badge_chips_wrap_to_second_row(self) -> None:
        geo = _cover_geometry(
            "Tiêu đề",
            subtitle="hook",
            badge_labels=["A" * 32, "B" * 32, "C" * 32, "D" * 32],
        )
        rows = {round(c[1], 3) for c in geo["badges"]}
        self.assertGreaterEqual(len(rows), 2)
        for _, by, bw, bh in geo["badges"]:
            self.assertLessEqual(by + bh, le.FOOTER_TOP_IN)
            self.assertLessEqual(bw, 3.10 + 1e-6)


class CoverHeroSlideTests(unittest.TestCase):
    def setUp(self) -> None:
        pres = _compute(_arc_payload())
        self.slide = pres.slides[0]

    def test_slide_uses_cover_layout_and_subtitle(self) -> None:
        self.assertEqual(self.slide.layout_type, LayoutType.COVER_HERO)
        self.assertTrue(self.slide.subtitle)

    def test_has_giant_title_subtitle_and_badges(self) -> None:
        title = next(e for e in self.slide.elements if e.morph_id == "title_slide")
        self.assertGreaterEqual(title.font_size, 44)
        self.assertLessEqual(title.font_size, 52)

        subtitle = next(e for e in self.slide.elements if e.morph_id == "cover_subtitle")
        self.assertEqual(subtitle.type, "heading")
        self.assertEqual(subtitle.font_size, le.COVER_SUB_PT)

        rule = next(e for e in self.slide.elements if e.morph_id == "title_rule")
        self.assertEqual(rule.type, "accent")

        chips = [e for e in self.slide.elements if e.morph_id in ("card_1", "card_2")]
        self.assertEqual(len(chips), 2)
        for chip in chips:
            self.assertEqual(chip.type, "badge")
            self.assertLessEqual(len(chip.content.split()), 3)

    def test_no_card_elements_on_cover(self) -> None:
        self.assertEqual(_elems(self.slide, "card"), [])


# ==============================================================================
# 4. GEOMETRY — BIG_STAT_CALLOUT & ASYMMETRIC_GRID
# ==============================================================================

class BigStatCalloutTests(unittest.TestCase):
    def test_number_font_is_capped_to_prevent_wrapping(self) -> None:
        for value in ("92 bar", "465 °C", "+35%", "90%", "99.86%", "2.2%"):
            size = _stat_number_font_size(value)
            self.assertGreaterEqual(size, 56, value)
            self.assertLessEqual(size, 64, value)
        self.assertEqual(_stat_number_font_size("92"), 64)

    def test_long_number_shrinks_but_stays_readable(self) -> None:
        """Chuỗi dài được hạ cỡ chữ; PPTX builder sẽ đặt word_wrap=False."""
        box_w = 5.21
        short_size = _stat_number_font_size("92 bar", box_w)
        for value in ("1.989×10³⁰ kg", "Tỷ lệ chuyển đổi của khách hàng doanh nghiệp"):
            size = _stat_number_font_size(value, box_w)
            self.assertGreaterEqual(size, 40)
            self.assertLess(size, short_size)
            self.assertLessEqual(size, 64)


    def test_stat_box_left_and_explanations_right(self) -> None:
        coords = _calculate_cards_geometry_inches(LayoutType.BIG_STAT_CALLOUT, 3)
        self.assertEqual(len(coords), 3)
        sx, sy, sw, sh = coords[0]
        self.assertEqual(sx, le.MX_IN)
        self.assertEqual(sh, le.CARDS_H_IN)

        usable = SLIDE_W - 2 * le.MX_IN
        ratio = sw / (usable - le.GAP_IN)
        self.assertAlmostEqual(ratio, le.STAT_COL_RATIO, places=2)

        for x, y, w, h in coords[1:]:
            self.assertGreater(x, sx + sw, "khối giải thích phải nằm bên phải con số")
            self.assertLessEqual(x + w, SLIDE_W - le.MX_IN + 1e-6)
            self.assertLessEqual(y + h, le.CARDS_TOP_IN + le.CARDS_H_IN + 1e-6)
        # hai thẻ giải thích xếp dọc, không chồng lấn
        self.assertGreater(coords[2][1], coords[1][1] + coords[1][3])

    def test_slide_renders_number_card_with_huge_font(self) -> None:
        slide = _compute(_arc_payload()).slides[1]
        self.assertEqual(slide.layout_type, LayoutType.BIG_STAT_CALLOUT)
        cards = _elems(slide, "card")
        self.assertEqual(len(cards), 2)
        self.assertGreaterEqual(cards[0].font_size, 56)
        self.assertLessEqual(cards[0].font_size, 64)
        self.assertEqual(cards[0].content, "92 bar")
        self.assertLessEqual(cards[1].font_size, 20)
        self.assertGreater(cards[1].x, cards[0].x + cards[0].width)


class AsymmetricGridTests(unittest.TestCase):
    def test_60_40_column_split(self) -> None:
        coords = _calculate_cards_geometry_inches(LayoutType.ASYMMETRIC_GRID, 3)
        self.assertEqual(len(coords), 3)
        hero_w = coords[0][2]
        side_w = coords[1][2]
        total = hero_w + side_w
        self.assertAlmostEqual(hero_w / total, 0.60, places=2)
        self.assertAlmostEqual(side_w / total, 0.40, places=2)

    def test_hero_is_full_height_and_side_cards_stack(self) -> None:
        coords = _calculate_cards_geometry_inches(LayoutType.ASYMMETRIC_GRID, 3)
        self.assertAlmostEqual(coords[0][3], le.CARDS_H_IN, places=3)
        self.assertEqual(coords[1][1], coords[0][1])
        self.assertGreater(coords[2][1], coords[1][1] + coords[1][3])
        self.assertEqual(coords[1][0], coords[2][0])
        self.assertLessEqual(coords[2][1] + coords[2][3], le.CARDS_TOP_IN + le.CARDS_H_IN + 1e-6)

    def test_two_cards_degrades_to_two_columns(self) -> None:
        coords = _calculate_cards_geometry_inches(LayoutType.ASYMMETRIC_GRID, 2)
        self.assertEqual(len(coords), 2)
        self.assertAlmostEqual(coords[0][3], le.CARDS_H_IN, places=3)
        self.assertAlmostEqual(coords[1][3], le.CARDS_H_IN, places=3)

    def test_slide_renders_hero_larger_than_side_cards(self) -> None:
        slide = _compute(_arc_payload()).slides[2]
        self.assertEqual(slide.layout_type, LayoutType.ASYMMETRIC_GRID)
        cards = _elems(slide, "card")
        self.assertEqual(len(cards), 3)
        self.assertGreater(cards[0].font_size, cards[1].font_size)
        self.assertGreater(cards[0].width, cards[1].width)


# ==============================================================================
# 5. GEOMETRY — CONCLUSION_SUMMARY
# ==============================================================================

class ConclusionSummaryTests(unittest.TestCase):
    def test_frame_is_centred_and_rows_fit_inside(self) -> None:
        fx, fy, fw, fh = _conclusion_frame_geometry()
        self.assertAlmostEqual(fx, SLIDE_W - fx - fw, places=3)  # căn giữa theo chiều ngang
        self.assertLess(fy, le.FOOTER_TOP_IN)

        rows = _conclusion_rows_geometry(3)
        self.assertEqual(len(rows), 3)
        for rx, ry, rw, rh in rows:
            self.assertGreaterEqual(rx, fx)
            self.assertLessEqual(rx + rw, fx + fw + 1e-6)
            self.assertGreaterEqual(ry, fy)
            self.assertLessEqual(ry + rh, fy + fh + 1e-6)
        for upper, lower in zip(rows, rows[1:]):
            self.assertGreater(lower[1], upper[1] + upper[3])

    def test_rows_stack_never_exceed_frame(self) -> None:
        fx, fy, fw, fh = _conclusion_frame_geometry()
        for n in (1, 2, 3):
            for ry, rh in [(r[1], r[3]) for r in _conclusion_rows_geometry(n)]:
                self.assertGreaterEqual(ry, fy)
                self.assertLessEqual(ry + rh, fy + fh + 1e-6, f"n={n}")

    def test_cta_sits_below_three_thin_rows_and_above_footer(self) -> None:
        slide = _compute(_arc_payload()).slides[4]
        self.assertEqual(slide.layout_type, LayoutType.CONCLUSION_SUMMARY)

        self.assertIsNone(next((e for e in slide.elements if e.morph_id == "summary_frame"), None))

        takeaways = [e for e in slide.elements if (e.morph_id or "").startswith("card_")
                     and e.type == "card" and e.step is not None]
        self.assertEqual(len(takeaways), 3)
        for row in takeaways:
            self.assertLessEqual(row.height, le.SUM_ROW_MAX_H_IN * 100 / SLIDE_H + 1e-6)
            self.assertTrue(row.border_color)
        for upper, lower in zip(takeaways, takeaways[1:]):
            self.assertGreater(lower.y, upper.y + upper.height)

        cta = next(e for e in slide.elements if e.morph_id == "call_to_action")
        self.assertEqual(cta.type, "card")
        self.assertGreater(cta.y, takeaways[-1].y + takeaways[-1].height)
        self.assertAlmostEqual(cta.height, le.SUM_CTA_H_IN * 100 / SLIDE_H, places=2)
        self.assertLessEqual(cta.y + cta.height, le.FOOTER_TOP_IN * 100 / SLIDE_H + 1e-6)
        self.assertIn("Theo dõi", cta.content)

    def test_three_cards_still_renders_without_cta(self) -> None:
        payload = _arc_payload()
        payload["slides"][4]["cards"] = payload["slides"][4]["cards"][:3]
        slide = _compute(payload).slides[4]
        self.assertEqual(slide.layout_type, LayoutType.CONCLUSION_SUMMARY)
        self.assertIsNone(next((e for e in slide.elements if e.morph_id == "call_to_action"), None))
        self.assertEqual(len(_elems(slide, "card")), 3)




# ==============================================================================
# 5b. AUTO STYLE RESOLUTION
# ==============================================================================

class AutoStyleResolutionTests(unittest.TestCase):
    def test_generate_request_accepts_auto_and_legacy_minimal_alias(self) -> None:
        self.assertEqual(GenerateRequest(prompt="báo cáo y học", style="AUTO").style, "AUTO")
        self.assertEqual(GenerateRequest(prompt="deck gọn", style="Minimal").style, "Minimalist")

    def test_auto_uses_llm_top_level_style_when_present(self) -> None:
        payload = _arc_payload()
        payload["style"] = "Academic"
        pres = _compute(payload, style="AUTO")
        self.assertEqual(pres.style, "Academic")

    def test_auto_falls_back_to_prompt_context_when_llm_style_missing(self) -> None:
        pres = _compute(_arc_payload(), style="AUTO")
        self.assertEqual(pres.style, "Futuristic")

# ==============================================================================
# 6. END-TO-END CONTRACT (schema ↔ frontend preview ↔ PPTX)
# ==============================================================================

class PresentationContractTests(unittest.TestCase):
    def setUp(self) -> None:
        self.pres = _compute(_arc_payload())

    def test_all_elements_inside_canvas_and_bounded(self) -> None:
        for slide in self.pres.slides:
            self.assertLessEqual(len(slide.elements), MAX_ELEMENTS_PER_SLIDE)
            for e in slide.elements:
                self.assertGreaterEqual(e.x, 0.0, e.id)
                self.assertGreaterEqual(e.y, 0.0, e.id)
                self.assertLessEqual(e.x + e.width, 100.01, e.id)
                self.assertLessEqual(e.y + e.height, 100.01, e.id)

    def test_cards_never_overlap_within_a_slide(self) -> None:
        for slide in self.pres.slides:
            cards = _elems(slide, "card")
            badges = _elems(slide, "badge")
            for group in (cards, badges):
                for i, a in enumerate(group):
                    for b in group[i + 1:]:
                        self.assertFalse(_overlap(a, b), f"{slide.slide_number}: {a.id} ∩ {b.id}")

    def test_json_round_trip_matches_frontend_contract(self) -> None:
        """JSON gửi cho frontend phải serialize được và validate lại được."""
        dumped = json.loads(self.pres.model_dump_json())
        self.assertEqual(dumped["total_slides"], 5)
        self.assertEqual(
            [s["layout_type"] for s in dumped["slides"]],
            ["COVER_HERO", "BIG_STAT_CALLOUT", "ASYMMETRIC_GRID",
             "TIMELINE_STEPS", "CONCLUSION_SUMMARY"],
        )
        for slide in dumped["slides"]:
            for elem in slide["elements"]:
                for key in ("id", "type", "content", "x", "y", "width", "height", "morph_id"):
                    self.assertIn(key, elem)
                self.assertIsInstance(elem["x"], (int, float))
        PresentationResponse.model_validate(dumped)

    def test_every_style_builds_without_geometry_error(self) -> None:
        for style in ("Futuristic", "Minimal", "Minimalist", "Corporate", "CorporateMinimalist", "Creative", "Academic"):
            pres = _compute(_arc_payload(), style=style)
            self.assertEqual(pres.slides[0].layout_type, LayoutType.COVER_HERO)
            self.assertEqual(pres.slides[-1].layout_type, LayoutType.CONCLUSION_SUMMARY)

    def test_pptx_is_built_from_the_new_layouts(self) -> None:
        from app.services.generator_service import build_pptx_from_schema
        from pptx import Presentation

        buffer = build_pptx_from_schema(self.pres)
        data = buffer.getvalue()
        self.assertGreater(len(data), 10_000)

        prs = Presentation(io.BytesIO(data))
        self.assertEqual(len(prs.slides.__iter__.__self__._sldIdLst), 5)
        shapes_per_slide = [len(list(s.shapes)) for s in prs.slides]
        for count in shapes_per_slide:
            self.assertGreaterEqual(count, 5)

    def test_morph_ids_are_stable_across_slides(self) -> None:
        morph_ids = [
            {e.morph_id for e in s.elements} for s in self.pres.slides
        ]
        for ids in morph_ids:
            self.assertIn("badge_header", ids)
            self.assertIn("footer_topic", ids)


# ==============================================================================
# 7. SYSTEM PROMPT & FEW-SHOT EXAMPLES
# ==============================================================================

class SystemPromptArcTests(unittest.TestCase):
    def test_prompt_mandates_the_three_part_arc(self) -> None:
        for needle in (
            "AUTO STYLE SELECTION", "top-level `style`",
            "MỞ BÀI", "THÂN BÀI", "KẾT BÀI",
            "CẤU TRÚC 3 PHẦN BẮT BUỘC",
            'layout_type BẮT BUỘC = "COVER_HERO"',
            'layout_type BẮT BUỘC = "CONCLUSION_SUMMARY"',
            "BIG_STAT_CALLOUT", "ASYMMETRIC_GRID", "TIMELINE_STEPS", "CARDS_ROW",
            "CALL TO ACTION",
            "TUYỆT ĐỐI KHÔNG dùng cùng một `layout_type` cho 2 slide LIÊN TIẾP",
        ):
            self.assertIn(needle, SYSTEM_PROMPT)

    def test_prompt_documents_card_roles_per_layout(self) -> None:
        for needle in (
            "CÂU HOOK (subtitle)", "BADGE", "CON SỐ / TỪ KHÓA KHỔNG LỒ",
            "HERO", "3 ĐIỂM CỐT LÕI", "ĐÚNG 4 thẻ", "ĐÚNG 3 thẻ",
        ):
            self.assertIn(needle, SYSTEM_PROMPT)

    def test_prompt_retires_legacy_layouts(self) -> None:
        self.assertIn("SPLIT_HERO`, `STAT_GRID`, `GRID_2X2` KHÔNG được dùng nữa", SYSTEM_PROMPT)

    def test_prompt_still_enforces_dynamic_persona(self) -> None:
        for needle in ("NHẬP VAI", "DYNAMIC ROLEPLAY", "DOMAIN VOCABULARY", "STRICT BAN"):
            self.assertIn(needle, SYSTEM_PROMPT)

    def test_few_shot_examples_are_schema_valid(self) -> None:
        """Ví dụ few-shot PHẢI khớp Pydantic Schema — LLM sẽ bắt chước chúng."""
        for blob in (_FEW_SHOT_EXAMPLE_1, _FEW_SHOT_EXAMPLE_2):
            payload = json.loads(blob)
            LLMOutput.model_validate(payload)

    def test_few_shot_example_1_shows_the_full_arc(self) -> None:
        payload = json.loads(_FEW_SHOT_EXAMPLE_1)
        layouts = [s["layout_type"] for s in payload["slides"]]
        self.assertEqual(layouts, ["COVER_HERO", "BIG_STAT_CALLOUT", "CONCLUSION_SUMMARY"])
        self.assertEqual(len(payload["slides"][0]["cards"]), 3)
        self.assertEqual(len(payload["slides"][1]["cards"]), 2)
        self.assertEqual(len(payload["slides"][2]["cards"]), 4)

    def test_few_shot_example_2_is_a_body_layout(self) -> None:
        payload = json.loads(_FEW_SHOT_EXAMPLE_2)
        self.assertEqual(payload["slides"][0]["layout_type"], "TIMELINE_STEPS")
        self.assertEqual(len(payload["slides"][0]["cards"]), 4)

    def test_user_prompt_states_arc_for_each_slide_count(self) -> None:
        one = _build_user_prompt(GenerateRequest(prompt="hệ mặt trời", num_slides=1))
        self.assertIn("chỉ 1 slide", one)
        self.assertIn("COVER_HERO", one)

        two = _build_user_prompt(GenerateRequest(prompt="hệ mặt trời", num_slides=2))
        self.assertIn("slide 2 BẮT BUỘC 'CONCLUSION_SUMMARY'", two)

        five = _build_user_prompt(GenerateRequest(prompt="hệ mặt trời", num_slides=5))
        self.assertIn("slide 1 = 'COVER_HERO'", five)
        self.assertIn("slide 2..4", five)
        self.assertIn("slide 5 = 'CONCLUSION_SUMMARY'", five)
        self.assertIn("BIG_STAT_CALLOUT", five)


if __name__ == "__main__":
    unittest.main(verbosity=2)
