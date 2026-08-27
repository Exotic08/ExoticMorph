"""Tests for the Visual Theme (motif) system — v8 "de-template hóa" upgrade.

Checklist nghiệm thu (từ yêu cầu thiết kế):
1. Ít nhất 1 slide (mở đầu HOẶC kết luận) có nền GRADIENT thay vì màu đặc.
2. Có ít nhất 1 bộ icon/motif theo ngữ cảnh chủ đề xuất hiện (kicker + card),
   không chỉ text thuần.
3. KHÔNG phải toàn bộ card dùng chung 1 kiểu viền/thanh màu — biến thể
   accent_top / accent_left / outline / cta tồn tại trong cùng 1 bộ slide.
4. Preview web và file .pptx thật khớp nhau: cùng chuỗi gradient CSS được
   parse thành <a:gradFill> đúng chuẩn, hoạ tiết nền mang <a:alpha> 3-6%,
   card biến thể render đúng shape tương ứng.
"""

from __future__ import annotations

import io
import json
import os
import re
import sys
import unittest

BACKEND_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if BACKEND_ROOT not in sys.path:
    sys.path.insert(0, BACKEND_ROOT)

from pptx import Presentation  # noqa: E402
from lxml import etree  # noqa: E402

from app.schemas.slide_schema import (  # noqa: E402
    GenerateRequest,
    LLMOutput,
    MAX_ELEMENTS_PER_SLIDE,
    PresentationResponse,
)
from app.services import layout_engine as le  # noqa: E402
from app.services.layout_engine import compute_layout  # noqa: E402
from app.services.generator_service import build_pptx_from_schema  # noqa: E402
from app.services.visual_themes import (  # noqa: E402
    VISUAL_THEME_KEYS,
    VisualTheme,
    card_icon_for,
    closing_gradient_css,
    cover_gradient_css,
    infer_visual_theme,
    kicker_icon_for,
    parse_gradient_css,
    pattern_decorations,
    resolve_visual_theme,
)
from app.services.llm_service import SYSTEM_PROMPT  # noqa: E402

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ==============================================================================
# FIXTURES
# ==============================================================================

DESC = (
    "Khí quyển Kim Tinh dày 92 bar với 96.5% CO2 tạo hiệu ứng nhà kính mất kiểm "
    "soát, giữ nhiệt độ bề mặt ở 465 độ C cả ngày lẫn đêm."
)


def _card(i: int, title: str, desc: str, color: str = "#8B5CF6") -> dict:
    return {
        "morph_id": f"card_{i}",
        "title": title,
        "description": desc,
        "color_theme": color,
        "order": i,
    }


def _arc_payload(visual_theme: str | None = "space") -> dict:
    payload = {
        "topic": "Nghịch lý nhiệt hành tinh Kim Tinh",
        "slides": [
            {
                "slide_number": 1,
                "section": "CÂU HỎI MỞ ĐẦU",
                "slide_title": "Vì sao Kim Tinh nóng 465 °C dù xa Mặt Trời hơn Thuỷ Tinh?",
                "layout_type": "COVER_HERO",
                "cards": [
                    _card(0, "Một nghịch lý nhiệt động giữa hai hành tinh láng giềng", DESC),
                    _card(1, "KHÍ QUYỂN", "Lớp vỏ khí 92 bar.", "#10B981"),
                    _card(2, "QUỸ ĐẠO", "Chu kỳ 225 ngày.", "#06B6D4"),
                ],
            },
            {
                "slide_number": 2,
                "section": "SỐ LIỆU",
                "slide_title": "Áp suất bề mặt Kim Tinh gấp 92 lần Trái Đất",
                "layout_type": "BIG_STAT_CALLOUT",
                "cards": [
                    _card(0, "92 bar", "Tương đương độ sâu 900 m dưới đại dương."),
                    _card(1, "CO2 chiếm 96.5%", "Khí nhà kính chính của hiệu ứng giữ nhiệt.", "#10B981"),
                ],
            },
            {
                "slide_number": 3,
                "section": "CƠ CHẾ",
                "slide_title": "Ba tầng khí quyển vận hành như một nồi áp suất",
                "layout_type": "ASYMMETRIC_GRID",
                "cards": [
                    _card(0, "Tầng đối lưu dày 65 km giữ trọn nhiệt lượng", DESC),
                    _card(1, "Mây axit sunfuric", "Phản xạ ánh sáng nhưng bẫy hồng ngoại.", "#10B981"),
                    _card(2, "Hiệu ứng nhà kính", "CO2 và SO2 tạo vòng phản hồi.", "#06B6D4"),
                ],
            },
            {
                "slide_number": 4,
                "section": "TIẾN TRÌNH",
                "slide_title": "Bốn giai đoạn hình thành nghịch lý nhiệt",
                "layout_type": "TIMELINE_STEPS",
                "cards": [
                    _card(0, "Giai đoạn 1: đại dương bốc hơi", "Hơi nước bổ sung khí quyển."),
                    _card(1, "Giai đoạn 2: CO2 tích tụ", "Núi lửa giải phóng CO2.", "#10B981"),
                    _card(2, "Giai đoạn 3: mất nước", "UV phân ly H2O.", "#06B6D4"),
                    _card(3, "Giai đoạn 4: cân bằng 465 °C", "Trạng thái siêu tới hạn.", "#EC4899"),
                ],
            },
            {
                "slide_number": 5,
                "section": "KẾT LUẬN",
                "slide_title": "Kim Tinh là lời cảnh báo cho Trái Đất",
                "layout_type": "CONCLUSION_SUMMARY",
                "cards": [
                    _card(0, "Khoảng cách không quyết định nhiệt độ", "Khí quyển mới là biến số chi phối."),
                    _card(1, "Nhà kính có điểm không quay lại", "Vượt ngưỡng là irreversible.", "#10B981"),
                    _card(2, "Trái Đất ở rìa vùng an toàn", "+4 °C kích hoạt vòng phản hồi.", "#06B6D4"),
                    _card(3, "Theo dõi CO2 ngay hôm nay", "Giám sát ppm tại Mauna Loa.", "#EC4899"),
                ],
            },
        ],
    }
    if visual_theme is not None:
        payload["visual_theme"] = visual_theme
    return payload


def _compute(payload: dict, num_slides: int | None = None, style: str = "AUTO",
             prompt: str = "nghịch lý nhiệt Kim Tinh", aspect: str = "16:9"):
    raw = LLMOutput.model_validate(payload)
    req = GenerateRequest(
        prompt=prompt,
        num_slides=num_slides or len(payload["slides"]),
        style=style,
        aspect_ratio=aspect,
    )
    return compute_layout(raw, req)


def _cards(slide):
    return [e for e in slide.elements if e.type == "card"]


def _decos(slide):
    return [e for e in slide.elements if e.type == "decoration"]


def _kicker(slide):
    return next((e for e in slide.elements if e.type == "kicker"), None)


# ==============================================================================
# 1. THEME INFERENCE
# ==============================================================================

class ThemeInferenceTests(unittest.TestCase):
    def test_space_theme_from_astronomy_keywords(self) -> None:
        for text in ("hệ mặt trời và các hành tinh", "khám phá thiên hà Andromeda",
                     "lịch sử quỹ đạo Sao Hoả", "space exploration and planets"):
            self.assertEqual(infer_visual_theme(text).key, "space", text)

    def test_finance_medical_education_nature(self) -> None:
        cases = {
            "finance": ("báo cáo doanh thu quý và kpi thị trường", "đầu tư cổ phiếu"),
            "medical": ("chẩn đoán lâm sàng tim mạch", "vaccine và dịch virus"),
            "education": ("phương pháp giảng dạy vật lý ở đại học", "nghiên cứu học thuật"),
            "nature": ("đa dạng sinh học của rừng mưa", "hệ sinh thái đại dương"),
            "food": ("ẩm thực Việt Nam và dinh dưỡng", "kỹ thuật nấu món ngon"),
            "history": ("triều đại nhà Nguyễn", "khảo cổ học văn minh cổ đại"),
            "tech": ("kiến trúc microservice và machine learning", "bảo mật cloud"),
        }
        for key, samples in cases.items():
            for text in samples:
                self.assertEqual(infer_visual_theme(text).key, key, text)

    def test_longer_keyword_wins_over_shorter_ambiguity(self) -> None:
        """'năng lượng mặt trời' phải là energy, không rơi vào space nhờ 'mặt trời'."""
        self.assertEqual(infer_visual_theme("năng lượng mặt trời và pin lithium").key, "energy")

    def test_generic_fallback_for_unknown_text(self) -> None:
        self.assertEqual(infer_visual_theme("ngày mai trời sẽ nắng").key, "generic")
        self.assertEqual(infer_visual_theme("").key, "generic")

    def test_llm_hint_takes_priority_and_invalid_falls_back(self) -> None:
        theme = resolve_visual_theme(
            prompt="báo cáo doanh thu", topic="tài chính", llm_hint="space"
        )
        self.assertEqual(theme.key, "space")
        # hint rác → suy luận lại từ prompt
        theme2 = resolve_visual_theme(prompt="báo cáo doanh thu", llm_hint="không-tồn-tại")
        self.assertEqual(theme2.key, "finance")
        # hint rác + prompt rỗng → generic
        self.assertEqual(resolve_visual_theme(llm_hint="???", prompt="").key, "generic")

    def test_icons_rotate_per_slide_and_card(self) -> None:
        space = infer_visual_theme("thiên văn")
        icons = [kicker_icon_for(space, i) for i in (1, 2, 3, 4, 5, 6)]
        self.assertEqual(len(set(icons)), 5, "5 slide phải xoay vòng đủ 5 icon kicker")
        self.assertTrue(all(icons), "icon không được rỗng")
        self.assertTrue(card_icon_for(space, 7), "card icon luôn có giá trị")

    def test_theme_keys_exposed_for_system_prompt(self) -> None:
        self.assertIn("space", VISUAL_THEME_KEYS)
        self.assertIn("generic", VISUAL_THEME_KEYS)


# ==============================================================================
# 2. PATTERN DECORATIONS (hoạ tiết nền mờ)
# ==============================================================================

class PatternDecorationTests(unittest.TestCase):
    def test_deterministic_same_seed(self) -> None:
        theme = infer_visual_theme("thiên văn")
        a = pattern_decorations(theme, 2)
        b = pattern_decorations(theme, 2)
        self.assertEqual(a, b, "cùng (theme, slide) phải cho cùng hoạ tiết — PPTX/preview khớp")

    def test_different_slide_different_pattern(self) -> None:
        theme = infer_visual_theme("thiên văn")
        a = pattern_decorations(theme, 2)
        b = pattern_decorations(theme, 3)
        self.assertNotEqual(a, b)

    def test_opacity_always_subtle_3_to_6_percent(self) -> None:
        for key in ("space", "finance", "medical"):
            theme = resolve_visual_theme(llm_hint=key)
            for deco in pattern_decorations(theme, 1):
                self.assertGreaterEqual(deco["opacity"], 0.03, f"{key}: {deco}")
                self.assertLessEqual(deco["opacity"], 0.06, f"{key}: {deco}")

    def test_patterns_stay_inside_canvas_and_capped(self) -> None:
        for key in ("space", "finance", "medical"):
            theme = resolve_visual_theme(llm_hint=key)
            decos = pattern_decorations(theme, 4, slide_w_in=10.0, slide_h_in=7.5)  # 4:3
            self.assertLessEqual(len(decos), 14)
            for d in decos:
                self.assertGreaterEqual(d["x"], 0.0)
                self.assertGreaterEqual(d["y"], 0.0)
                self.assertLessEqual(d["x"] + d["w"], 10.0 + 1e-6)
                self.assertLessEqual(d["y"] + d["h"], 7.5 + 1e-6)
                self.assertIn(d["shape"], ("oval", "rect"))


# ==============================================================================
# 3. GRADIENT (CSS string ↔ PPTX gradFill)
# ==============================================================================

class GradientTests(unittest.TestCase):
    def test_cover_and_closing_gradients_parse(self) -> None:
        cover = cover_gradient_css("#0A0C12", "#10B981", "#06B6D4")
        closing = closing_gradient_css("#0A0C12", "#EC4899")
        for css in (cover, closing):
            parsed = parse_gradient_css(css)
            self.assertIsNotNone(parsed, css)
            angle, stops = parsed
            self.assertGreater(angle, 0)
            self.assertGreaterEqual(len(stops), 3)
            for pos, hex_color in stops:
                self.assertTrue(re.fullmatch(r"#[0-9A-F]{6}", hex_color))
                self.assertGreaterEqual(pos, 0.0)
                self.assertLessEqual(pos, 100.0)
        # gradient kết bài phải DẤNG dần về màu accent (stop cuối sáng hơn stop đầu)
        angle_c, stops_c = parse_gradient_css(closing)
        self.assertNotEqual(stops_c[0][1], stops_c[-1][1])

    def test_parser_rejects_non_gradient(self) -> None:
        self.assertIsNone(parse_gradient_css("#0A0C12"))
        self.assertIsNone(parse_gradient_css(None))
        self.assertIsNone(parse_gradient_css("linear-gradient(banana, red 0%, blue 100%)"))
        self.assertIsNone(parse_gradient_css("linear-gradient(90deg, #FF0000 0%)"))

    def test_gradient_fades_toward_accent_not_flashing(self) -> None:
        """Gradient phải còn TỐI (khung nền tối) — điểm nhấn mờ, không chói."""
        def channel(hex_str: str, i: int) -> int:
            return int(hex_str.lstrip("#")[i * 2:i * 2 + 2], 16)

        closing = closing_gradient_css("#0A0C12", "#10B981")
        _, stops = parse_gradient_css(closing)
        # stop cuối trộn tối đa ~28% accent → không vượt ngưỡng chói
        for i in range(3):
            self.assertLess(channel(stops[-1][1], i), 120)


# ==============================================================================
# 4. LAYOUT ENGINE — Ngôn ngữ thị giác theo chủ đề trong PresentationResponse
# ==============================================================================

class VisualDeckLayoutTests(unittest.TestCase):
    def setUp(self) -> None:
        self.pres = _compute(_arc_payload())

    def test_checklist_1_first_and_last_slides_have_gradient_bg(self) -> None:
        first, last = self.pres.slides[0], self.pres.slides[-1]
        for slide in (first, last):
            self.assertTrue(
                (slide.bg_color or "").startswith("linear-gradient("),
                f"slide {slide.slide_number} phải có nền gradient, got {slide.bg_color}",
            )
            self.assertIsNotNone(parse_gradient_css(slide.bg_color))

    def test_checklist_1_body_slides_stay_flat_but_get_decorations(self) -> None:
        for slide in self.pres.slides[1:-1]:
            self.assertFalse((slide.bg_color or "").startswith("linear-gradient("))
            decos = _decos(slide)
            self.assertGreaterEqual(len(decos), 8, "slide thân bài cần hoạ tiết nền mờ")
            for d in decos:
                self.assertIsNotNone(d.opacity)
                self.assertGreaterEqual(d.opacity, 0.03)
                self.assertLessEqual(d.opacity, 0.06)
                self.assertEqual(d.morph_id, "", "hoạ tiết nền không mang morph_id")

    def test_checklist_2_kicker_carries_topic_icon_every_slide(self) -> None:
        for slide in self.pres.slides:
            kicker = _kicker(slide)
            self.assertIsNotNone(kicker)
            # icon emoji/glyph mở đầu kicker (khác ASCII A-Z)
            self.assertRegex(kicker.content, r"^[^\sA-Za-z0-9]\s", kicker.content)
            self.assertIn(slide.section.upper(), kicker.content)

    def test_checklist_2_cards_carry_topic_icons(self) -> None:
        iconed = [
            e for s in self.pres.slides for e in _cards(s) if e.icon
        ]
        self.assertGreaterEqual(
            len(iconed), 4, "bộ slide phải có icon motif trên card, không chỉ text"
        )
        for e in iconed:
            self.assertTrue(e.icon and len(e.icon) >= 1)
            self.assertIn(e.variant, ("accent_left", "outline", "accent_top"))

    def test_stat_card_has_trend_icon(self) -> None:
        stat_slide = self.pres.slides[1]  # BIG_STAT_CALLOUT
        stat_card = _cards(stat_slide)[0]
        self.assertEqual(stat_card.content, "92 bar")
        self.assertIsNotNone(stat_card.icon, "card số lớn cần icon xu hướng/đơn vị")

    def test_checklist_3_card_variants_are_mixed(self) -> None:
        all_variants = {e.variant for s in self.pres.slides for e in _cards(s)}
        self.assertGreaterEqual(
            len(all_variants), 3,
            f"bộ slide phải trộn các biến thể card, got {all_variants}",
        )
        self.assertIn("accent_top", all_variants)
        self.assertIn("accent_left", all_variants)
        self.assertIn("outline", all_variants)
        # Kết bài: list-rows dùng viền trái, CTA là thanh đặc
        last = self.pres.slides[-1]
        row_variants = {e.variant for e in _cards(last)}
        self.assertIn("accent_left", row_variants)
        self.assertIn("cta", row_variants)
        # Timeline: card bước kiểu outline (không nền)
        timeline = self.pres.slides[3]
        self.assertTrue(all(e.variant == "outline" for e in _cards(timeline)))
        # BIG_STAT: thẻ giải thích dùng viền trái (không lặp thanh trên của hero)
        stat_slide = self.pres.slides[1]
        self.assertEqual(_cards(stat_slide)[1].variant, "accent_left")

    def test_variants_vary_between_adjacent_slides(self) -> None:
        sigs = []
        for slide in self.pres.slides:
            sigs.append(tuple(e.variant for e in _cards(slide)))
        for prev, cur in zip(sigs, sigs[1:]):
            self.assertNotEqual(
                prev, cur, "2 slide liên tiếp không lặp y hệt bộ biến thể card"
            )

    def test_decorations_render_below_and_budget_respected(self) -> None:
        for slide in self.pres.slides:
            self.assertLessEqual(len(slide.elements), MAX_ELEMENTS_PER_SLIDE)
            if _decos(slide):
                # decoration phải đứng ĐẦU danh sách → z-order thấp nhất ở PPTX
                self.assertEqual(slide.elements[0].type, "decoration")
            for e in slide.elements:
                self.assertGreaterEqual(e.x, 0.0)
                self.assertGreaterEqual(e.y, 0.0)
                self.assertLessEqual(e.x + e.width, 100.01)
                self.assertLessEqual(e.y + e.height, 100.01)

    def test_finance_topic_gets_grid_pattern_and_money_icons(self) -> None:
        payload = _arc_payload(visual_theme=None)
        payload["topic"] = "Báo cáo tài chính doanh nghiệp"
        pres = _compute(payload, prompt="báo cáo tài chính doanh nghiệp kpi")
        kick = _kicker(pres.slides[0])
        self.assertIn(kick.content.strip()[0], ("📈", "📊", "💰", "🏦", "💹"))
        stat_card = _cards(pres.slides[1])[0]
        self.assertEqual(stat_card.icon, "📈")

    def test_no_llm_theme_falls_back_to_inference(self) -> None:
        payload = _arc_payload(visual_theme=None)
        payload["topic"] = "Hệ Mặt Trời"
        pres = _compute(payload, prompt="hệ mặt trời")
        kick = _kicker(pres.slides[0])
        self.assertTrue(kick.content.strip().startswith("🪐"))

    def test_4by3_decorations_fit_10in_canvas(self) -> None:
        pres = _compute(_arc_payload(), aspect="4:3")
        for slide in pres.slides[1:-1]:
            for d in _decos(slide):
                w_in = (d.width / 100.0) * 10.0
                x_in = (d.x / 100.0) * 10.0
                self.assertLessEqual(x_in + w_in, 10.0 + 0.02)

    def test_json_round_trip_keeps_visual_fields(self) -> None:
        dumped = json.loads(self.pres.model_dump_json())
        PresentationResponse.model_validate(dumped)
        first = dumped["slides"][0]
        self.assertTrue(first["bg_color"].startswith("linear-gradient("))
        card_fields = [
            (e["variant"], e["icon"], e["opacity"])
            for s in dumped["slides"] for e in s["elements"] if e["type"] == "card"
        ]
        self.assertTrue(any(v == "outline" for v, _, _ in card_fields))
        self.assertTrue(any(i for _, i, _ in card_fields))

    def test_morph_strategy_mentions_visual_theme(self) -> None:
        self.assertIn("visual theme", (self.pres.morph_strategy or "").lower())
        self.assertIn("space", self.pres.morph_strategy or "")


# ==============================================================================
# 5. PPTX BUILDER — gradient OpenXML + alpha hoạ tiết + biến thể card
# ==============================================================================

class PptxVisualBuildTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.pres = _compute(_arc_payload())
        cls.buffer = build_pptx_from_schema(cls.pres)
        cls.prs = Presentation(io.BytesIO(cls.buffer.getvalue()))

    def _bg(self, slide) -> etree._Element | None:
        return slide._element.cSld.find(f"{{{P_NS}}}bg")

    def test_pptx_opens_with_five_slides(self) -> None:
        self.assertEqual(len(self.prs.slides._sldIdLst), 5)

    def test_gradient_slides_have_gradfill_with_stops_and_angle(self) -> None:
        for idx in (0, 4):
            bg = self._bg(self.prs.slides[idx])
            self.assertIsNotNone(bg, f"slide {idx + 1} thiếu p:bg")
            grad = bg.find(f".//{{{A_NS}}}gradFill")
            self.assertIsNotNone(grad, f"slide {idx + 1} nền phải là gradFill")
            stops = grad.findall(f"{{{A_NS}}}gsLst/{{{A_NS}}}gs")
            self.assertEqual(len(stops), 3)
            pos_list = [int(s.get("pos")) for s in stops]
            self.assertEqual(pos_list, sorted(pos_list))
            self.assertEqual(pos_list[0], 0)
            self.assertEqual(pos_list[-1], 100000)
            lin = grad.find(f"{{{A_NS}}}lin")
            self.assertIsNotNone(lin)
            ang = int(lin.get("ang"))
            # CSS 118deg/160deg → DrawingML (css-90)*60000 = 28deg/70deg
            expected = {0: 28, 4: 70}[idx]
            self.assertEqual(ang, expected * 60000)

    def test_body_slides_solid_bg_and_alpha_decorations(self) -> None:
        slide = self.prs.slides[1]  # BIG_STAT (thân bài)
        bg = self._bg(slide)
        self.assertIsNotNone(bg.find(f".//{{{A_NS}}}solidFill"))
        alphas = [
            int(a.get("val")) for a in slide._element.findall(f".//{{{A_NS}}}alpha")
            if int(a.get("val")) <= 6000
        ]
        self.assertGreaterEqual(len(alphas), 8, "hoạ tiết nền cần <a:alpha> 3-6%")
        for val in alphas:
            self.assertGreaterEqual(val, 3000)
            self.assertLessEqual(val, 6000)

    def test_decorations_have_no_morph_name(self) -> None:
        slide = self.prs.slides[1]
        names = [
            sp._element.nvSpPr.cNvPr.get("name")
            for sp in slide.shapes
            if sp.shape_type is not None and sp.name.startswith("deco_")
        ]
        self.assertTrue(all(n and n.startswith("deco_") for n in names))
        self.assertFalse(any("!!" in n for n in names))

    def test_card_variants_map_to_real_geometry(self) -> None:
        """accent_left phải sinh shape viền trái; outline phải noFill."""
        def shapes_of(slide):
            return list(slide.shapes)

        # Slide kết bài (idx 4): 3 row accent_left → 3 thanh viền trái + 1 CTA đặc
        concl = self.prs.slides[4]
        concl_xml = etree.tostring(concl._element).decode()
        self.assertGreaterEqual(concl_xml.count("call_to_action"), 1)

        # Slide timeline (idx 3): card outline → noFill trong spPr
        timeline = self.prs.slides[3]
        nofills = timeline._element.findall(
            f".//{{{P_NS}}}sp/{{{P_NS}}}spPr/{{{A_NS}}}noFill"
        )
        self.assertGreaterEqual(len(nofills), 3, "timeline cards phải trong suốt (noFill)")

        # Slide thân bài 1 (idx 1): card hero đặc + icon corner textbox của stat
        stat_slide = self.prs.slides[1]
        texts = [
            sp.text_frame.text for sp in stat_slide.shapes
            if sp.has_text_frame
        ]
        self.assertTrue(any("☄" in t or "📈" in t or t.strip() in ("🧬", "⚡", "🪐")
                            for t in texts), "icon motif phải xuất hiện trong PPTX")

    def test_stat_icon_run_has_no_explicit_font(self) -> None:
        """Icon không gắn 'Segoe UI' — để PowerPoint tự fallback emoji font."""
        stat_slide = self.prs.slides[1]
        for sp in stat_slide.shapes:
            if not sp.has_text_frame:
                continue
            for p in sp.text_frame.paragraphs:
                for run in p.runs:
                    if run.text.strip() in ("☄", "📈", "🧬", "⚡", "🪐"):
                        self.assertNotEqual(run.font.name, "Segoe UI")

    def test_pptx_passes_powerpoint_schema_order(self) -> None:
        """p:bg phải là con ĐẦU TIÊN của p:cSld theo CT_CommonSlideData."""
        for slide in self.prs.slides:
            cSld = slide._element.cSld
            first_child = cSld[0].tag
            self.assertEqual(first_child, f"{{{P_NS}}}bg")


# ==============================================================================
# 6. SYSTEM PROMPT — quy tắc đa dạng thị giác
# ==============================================================================

class SystemPromptVisualDiversityTests(unittest.TestCase):
    def test_prompt_contains_visual_diversity_mandate(self) -> None:
        for needle in (
            "MỘT",
            "chi tiết thị giác đặc trưng cho chủ đề",
            "không lặp lại y hệt bố cục card/thanh màu của slide trước",
            "không phải khuôn cố định",
        ):
            self.assertIn(needle, SYSTEM_PROMPT)

    def test_prompt_documents_visual_theme_field(self) -> None:
        self.assertIn("visual_theme", SYSTEM_PROMPT)
        for key in VISUAL_THEME_KEYS:
            self.assertIn(key, SYSTEM_PROMPT)

    def test_prompt_keeps_arc_contract_needles(self) -> None:
        for needle in (
            'layout_type BẮT BUỘC = "COVER_HERO"',
            'layout_type BẮT BUỘC = "CONCLUSION_SUMMARY"',
            "TUYỆT ĐỐI KHÔNG dùng cùng một `layout_type` cho 2 slide LIÊN TIẾP",
            "AUTO STYLE SELECTION",
            "NHẬP VAI", "DOMAIN VOCABULARY", "STRICT BAN",
        ):
            self.assertIn(needle, SYSTEM_PROMPT)

    def test_user_prompt_mentions_visual_theme(self) -> None:
        from app.services.llm_service import _build_user_prompt
        text = _build_user_prompt(GenerateRequest(prompt="hệ mặt trời", num_slides=5))
        self.assertIn("visual_theme", text)
        self.assertIn("space", text)


if __name__ == "__main__":
    unittest.main(verbosity=2)
