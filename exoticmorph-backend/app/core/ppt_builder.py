"""ExoticMorph PowerPoint Builder Module
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
Module dựng slide thuyết trình bằng python-pptx và tích hợp bộ xử lý XML Morph.
"""

from typing import Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .xml_helper import enable_morph_transition, set_shape_morph_id


class MorphPresentationBuilder:
  """Lớp tiện ích xây dựng bài trình chiếu PowerPoint tỉ lệ chuẩn 16:9 hỗ trợ

  toàn diện hiệu ứng chuyển cảnh Morph.
  """

  # Bảng màu chuẩn thiết kế Modern Dark & Vibrant
  COLOR_BG_DARK = RGBColor(14, 16, 23)  # #0E1017 (Dark Background)
  COLOR_SLIDE_BG = RGBColor(18, 20, 28)  # #12141C (Slide Background)
  COLOR_TEXT_WHITE = RGBColor(255, 255, 255)
  COLOR_TEXT_MUTED = RGBColor(148, 163, 184)  # #94A3B8 (Slate 400)
  COLOR_GREEN_START = RGBColor(16, 185, 129)  # #10B981 (Emerald 500)
  COLOR_ORANGE_END = RGBColor(249, 115, 22)  # #F97316 (Orange 500)
  COLOR_PURPLE_ACCENT = RGBColor(139, 92, 246)  # #8B5CF6 (Purple 500)
  COLOR_CYAN_ACCENT = RGBColor(6, 182, 212)  # #06B6D4 (Cyan 500)

  def __init__(
      self, width_inches: float = 13.333, height_inches: float = 7.5
  ) -> None:
    """Khởi tạo Presentation với kích thước 16:9 Widescreen (13.333 x 7.5

    inches).
    """
    self.prs = Presentation()
    self.prs.slide_width = Inches(width_inches)
    self.prs.slide_height = Inches(height_inches)
    self.blank_layout = self.prs.slide_layouts[6]  # Blank Slide Layout

  def _set_slide_background(
      self, slide, color: RGBColor = COLOR_SLIDE_BG
  ) -> None:
    """Tô màu nền toàn màn hình cho slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

  def add_blank_slide(self, bg_color: Optional[RGBColor] = None):
    """Thêm một slide trắng có tô màu nền hiện đại."""
    slide = self.prs.slides.add_slide(self.blank_layout)
    self._set_slide_background(slide, bg_color or self.COLOR_SLIDE_BG)
    return slide

  def add_header(
      self,
      slide,
      title: str,
      subtitle: str,
      badge: Optional[str] = "ExoticMorph AI Demo",
  ) -> None:
    """Thêm tiêu đề và mô tả trạng thái chuyển động ở góc trên slide."""
    # Text container
    left = Inches(0.8)
    top = Inches(0.6)
    width = Inches(11.7)
    height = Inches(1.2)

    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    # Badge + Title paragraph
    p1 = tf.paragraphs[0]
    p1.text = f"[{badge}] {title}"
    p1.font.bold = True
    p1.font.size = Pt(22)
    p1.font.color.rgb = self.COLOR_TEXT_WHITE
    p1.font.name = "Segoe UI"

    # Subtitle paragraph
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(12)
    p2.font.color.rgb = self.COLOR_TEXT_MUTED
    p2.font.name = "Segoe UI"
    p2.space_before = Pt(4)

  def save(self, output_path: str) -> str:
    """Lưu bài thuyết trình ra tệp .pptx."""
    self.prs.save(output_path)
    return output_path


def create_demo_morph_presentation(
    output_path: str = "demo_morph.pptx",
) -> str:
  """Xây dựng bài thuyết trình mẫu gồm 2 slide để kiểm chứng hiệu ứng Morph

  OpenXML:

  - **Slide 1:**
      - Vị trí: Góc trên bên trái (Left=1.0", Top=2.2").
      - Kích thước: 3.2" x 2.0".
      - Màu sắc: Xanh lá (Emerald Green - #10B981).
      - Nội dung chữ: "Slide 1: Start".
      - Morph ID: `!!DemoBox`.

  - **Slide 2:**
      - Vị trí: Di chuyển xuống góc dưới bên phải (Left=7.0", Top=3.2").
      - Kích thước: Gấp đôi (6.4" x 4.0").
      - Màu sắc: Đổi sang màu cam (Vibrant Orange - #F97316).
      - Nội dung chữ: "Slide 2: End".
      - Morph ID: `!!DemoBox` (Trùng ID để PowerPoint tự động nội suy
      chuyển động).
      - Kích hoạt Transition Morph: `<p:transition><p:morph/></p:transition>`.

  :param output_path: Đường dẫn tệp .pptx cần lưu.
  :return: Đường dẫn tệp đã lưu thành công.
  """
  builder = MorphPresentationBuilder()

  # ==========================================================================
  # 1. SLIDE 1: TRẠNG THÁI KHỞI ĐẦU (START STATE)
  # ==========================================================================
  slide1 = builder.add_blank_slide()

  builder.add_header(
      slide=slide1,
      title="Slide 1: Trạng Thái Khởi Đầu",
      subtitle=(
          "Hình chữ nhật nằm ở góc trên bên trái, màu xanh lá cây. Khi chuyển"
          " slide, PowerPoint sẽ đọc ID '!!DemoBox' để kích hoạt hiệu ứng"
          " Morph."
      ),
      badge="Slide 1 / 2",
  )

  # Tạo hình chữ nhật ở góc trên bên trái
  s1_left = Inches(1.0)
  s1_top = Inches(2.2)
  s1_width = Inches(3.2)
  s1_height = Inches(2.0)

  shape1 = slide1.shapes.add_shape(
      MSO_SHAPE.ROUNDED_RECTANGLE, s1_left, s1_top, s1_width, s1_height
  )

  # Cấu hình màu sắc Xanh lá
  shape1.fill.solid()
  shape1.fill.fore_color.rgb = MorphPresentationBuilder.COLOR_GREEN_START
  shape1.line.color.rgb = RGBColor(52, 211, 153)  # Viền xanh sáng
  shape1.line.width = Pt(1.5)

  # Cấu hình nội dung Text
  tf1 = shape1.text_frame
  tf1.word_wrap = True
  p1 = tf1.paragraphs[0]
  p1.text = "Slide 1: Start"
  p1.alignment = PP_ALIGN.CENTER
  p1.font.bold = True
  p1.font.size = Pt(20)
  p1.font.color.rgb = RGBColor(255, 255, 255)
  p1.font.name = "Segoe UI"

  # Thêm dòng chú thích nhỏ bên trong shape
  p1_sub = tf1.add_paragraph()
  p1_sub.text = "Position: (1.0\", 2.2\") | 3.2\" x 2.0\""
  p1_sub.alignment = PP_ALIGN.CENTER
  p1_sub.font.size = Pt(10)
  p1_sub.font.color.rgb = RGBColor(230, 255, 240)
  p1_sub.font.name = "Segoe UI"
  p1_sub.space_before = Pt(6)

  # GÁN MORPH ID CHO SHAPE 1 (Tạo tiền tố !!)
  morph_id_assigned_1 = set_shape_morph_id(shape1, "DemoBox")

  # Thêm một thẻ phụ minh họa Morph Anchor Info
  info1 = slide1.shapes.add_textbox(
      Inches(1.0), Inches(4.5), Inches(5.0), Inches(1.5)
  )
  itf1 = info1.text_frame
  ip1 = itf1.paragraphs[0]
  ip1.text = f"Morph Object Name: {morph_id_assigned_1}"
  ip1.font.size = Pt(11)
  ip1.font.color.rgb = MorphPresentationBuilder.COLOR_CYAN_ACCENT
  ip1.font.bold = True

  ip2 = itf1.add_paragraph()
  ip2.text = "• XML Node: <p:cNvPr name=\"!!DemoBox\"/>"
  ip2.font.size = Pt(10)
  ip2.font.color.rgb = MorphPresentationBuilder.COLOR_TEXT_MUTED

  # ==========================================================================
  # 2. SLIDE 2: TRẠNG THÁI ĐÍCH ĐẾN (END STATE & MORPH TRANSITION)
  # ==========================================================================
  slide2 = builder.add_blank_slide()

  builder.add_header(
      slide=slide2,
      title="Slide 2: Trạng Thái Biến Đổi Hoàn Tất (Morph End)",
      subtitle=(
          "Đối tượng '!!DemoBox' đã tự động di chuyển xuống góc dưới phải, mở"
          " rộng gấp đôi kích thước và chuyển màu cam mượt mà."
      ),
      badge="Slide 2 / 2",
  )

  # Tạo hình chữ nhật ở góc dưới bên phải với kích thước gấp đôi
  s2_left = Inches(6.0)
  s2_top = Inches(2.2)
  s2_width = Inches(6.4)  # Kích thước to gấp đôi chiều ngang
  s2_height = Inches(4.0)  # Kích thước to gấp đôi chiều dọc

  shape2 = slide2.shapes.add_shape(
      MSO_SHAPE.ROUNDED_RECTANGLE, s2_left, s2_top, s2_width, s2_height
  )

  # Cấu hình màu sắc Cam
  shape2.fill.solid()
  shape2.fill.fore_color.rgb = MorphPresentationBuilder.COLOR_ORANGE_END
  shape2.line.color.rgb = RGBColor(251, 146, 60)  # Viền cam sáng
  shape2.line.width = Pt(2.0)

  # Cấu hình nội dung Text
  tf2 = shape2.text_frame
  tf2.word_wrap = True
  p2 = tf2.paragraphs[0]
  p2.text = "Slide 2: End"
  p2.alignment = PP_ALIGN.CENTER
  p2.font.bold = True
  p2.font.size = Pt(28)
  p2.font.color.rgb = RGBColor(255, 255, 255)
  p2.font.name = "Segoe UI"

  # Thêm dòng chú thích bên trong shape 2
  p2_sub1 = tf2.add_paragraph()
  p2_sub1.text = "Position: (6.0\", 2.2\") | 6.4\" x 4.0\" (2x Size)"
  p2_sub1.alignment = PP_ALIGN.CENTER
  p2_sub1.font.size = Pt(13)
  p2_sub1.font.bold = True
  p2_sub1.font.color.rgb = RGBColor(255, 240, 230)
  p2_sub1.font.name = "Segoe UI"
  p2_sub1.space_before = Pt(8)

  p2_sub2 = tf2.add_paragraph()
  p2_sub2.text = (
      "Tự động chuyển đổi mượt mà bằng XML Morphing Engine của PowerPoint!"
  )
  p2_sub2.alignment = PP_ALIGN.CENTER
  p2_sub2.font.size = Pt(11)
  p2_sub2.font.color.rgb = RGBColor(255, 230, 200)
  p2_sub2.font.name = "Segoe UI"
  p2_sub2.space_before = Pt(4)

  # GÁN CÙNG MORPH ID 'DemoBox' CHO SHAPE 2
  morph_id_assigned_2 = set_shape_morph_id(shape2, "DemoBox")

  # KÍCH HOẠT HIỆU ỨNG MORPH TRANSITION CHO SLIDE 2
  enable_morph_transition(
      slide=slide2, duration_ms=1250, morph_option="byObject"
  )

  # Thêm thẻ giải thích kỹ thuật XML trên Slide 2
  info2 = slide2.shapes.add_textbox(
      Inches(0.8), Inches(2.2), Inches(4.8), Inches(4.0)
  )
  itf2 = info2.text_frame
  itf2.word_wrap = True

  tp1 = itf2.paragraphs[0]
  tp1.text = "Thông Số Kỹ Thuật OpenXML:"
  tp1.font.size = Pt(14)
  tp1.font.bold = True
  tp1.font.color.rgb = MorphPresentationBuilder.COLOR_PURPLE_ACCENT

  bullets = [
      ("Morph Matching:", f"Shape Name = '{morph_id_assigned_2}'"),
      (
          "Slide Transition:",
          '<p:transition spd="med"><p:morph option="byObject"/></p:transition>',
      ),
      ("Schema Conformance:", "Chèn trước thẻ <p:clrMapOvr> hợp lệ ISO/IEC"),
      (
          "Interpolation:",
          "Nội suy đồng thời Position, Size, Color Fill & Font Size",
      ),
  ]

  for label, desc in bullets:
    p = itf2.add_paragraph()
    p.text = f"• {label} {desc}"
    p.font.size = Pt(10.5)
    p.font.color.rgb = MorphPresentationBuilder.COLOR_TEXT_WHITE
    p.space_before = Pt(8)

  # Lưu tệp PowerPoint
  return builder.save(output_path)
