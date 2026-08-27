"""
PPTX Generator Service for ExoticMorph
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
Đọc dữ liệu có cấu trúc từ Layout Engine (PresentationResponse), chuyển đổi tọa
độ % sang Inches, dựng slide bằng python-pptx và can thiệp OpenXML để kích hoạt
hiệu ứng Morph.

Đợt nâng cấp (Executive Design):
- Kiểu chữ phân cấp rõ ràng: kicker in hoa có tracking, tiêu đề lớn đậm, thân
  chữ tương phản cao với khoảng cách dòng thoáng.
- Mỗi thẻ có thanh accent trên đỉnh (màu nhận diện) + nền tối sang trọng.
- Bố cục lộ trình (roadmap) có vòng tròn số thứ tự và đường nối.
- Footer cố định (chủ đề + số trang) trên mọi slide để đồng bộ bộ nhận diện.
"""

import io
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml

from app.schemas.slide_schema import PresentationResponse, Slide, SlideElement
from app.core.xml_helper import enable_morph_transition, set_shape_morph_id
from app.services.layout_engine import STYLE_PALETTES
from app.services.visual_themes import parse_gradient_css

# Thời lượng chuyển cảnh Morph (ms) — nhanh, dứt khoát như keynote chuyên nghiệp.
MORPH_DURATION_MS = 900
EMU_PER_INCH = 914400
CARD_TEXT_PAD_IN = 0.18   # 12.96pt — đồng bộ preview web và PPTX
CARD_TEXT_PAD_TOP_IN = 0.20  # 14.4pt
CARD_TEXT_PAD_BOTTOM_IN = 0.16  # 11.52pt
CARD_BORDER_PT = 0.75
CARD_BORDER_OPACITY = 0.55
# Biến thể card: viền trái dày của accent_left (≈4.5px @96dpi) — đồng bộ preview.
CARD_LEFT_BAR_IN = 0.062
# Viền mỏng của card outline (không nền) — hơi dày hơn viền card đặc để định hình.
OUTLINE_BORDER_PT = 1.0
OUTLINE_BORDER_OPACITY = 0.55


def hex_to_rgb(hex_str: Optional[str], default: RGBColor = RGBColor(255, 255, 255)) -> RGBColor:
    """Chuyển đổi chuỗi Hex (ví dụ '#8B5CF6' hoặc '8B5CF6') sang RGBColor của python-pptx.

    Chịu lỗi: chuỗi rỗng, sai độ dài, ký tự hex không hợp lệ -> trả màu default.
    """
    if not hex_str or not isinstance(hex_str, str):
        return default

    clean = hex_str.strip().lstrip("#")
    if len(clean) == 3:
        clean = "".join([c * 2 for c in clean])
    if len(clean) != 6:
        return default

    try:
        r = int(clean[0:2], 16)
        g = int(clean[2:4], 16)
        b = int(clean[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        return default


def _set_line_opacity(shape, opacity: float = CARD_BORDER_OPACITY) -> None:
    """Thiết lập opacity cho stroke bằng OpenXML (0..1).

    python-pptx chưa expose trực tiếp alpha của line color, nên cần chèn
    ``<a:alpha>`` vào ``a:srgbClr``. Nếu XML khác kỳ vọng thì bỏ qua an toàn.
    """
    try:
        opacity = max(0.0, min(1.0, opacity))
        ln = shape._element.spPr.get_or_add_ln()
        solid = ln.find(qn("a:solidFill"))
        if solid is None:
            return
        srgb = solid.find(qn("a:srgbClr"))
        if srgb is None:
            return
        alpha = srgb.find(qn("a:alpha"))
        if alpha is None:
            alpha = OxmlElement("a:alpha")
            srgb.append(alpha)
        alpha.set("val", str(int(opacity * 100000)))
    except Exception:  # noqa: BLE001
        pass


def _set_fill_opacity(shape, opacity: float) -> None:
    """Thiết lập opacity cho FILL của shape bằng OpenXML (0..1).

    Dùng cho hoạ tiết nền (decoration) — opacity 3-6% tạo chiều sâu mà không
    rối mắt. Cơ chế giống ``_set_line_opacity`` nhưng trên ``a:solidFill``
    của thân shape.
    """
    try:
        opacity = max(0.0, min(1.0, opacity))
        spPr = shape._element.spPr
        solid = spPr.find(qn("a:solidFill"))
        if solid is None:
            return
        srgb = solid.find(qn("a:srgbClr"))
        if srgb is None:
            return
        alpha = srgb.find(qn("a:alpha"))
        if alpha is None:
            alpha = OxmlElement("a:alpha")
            srgb.append(alpha)
        alpha.set("val", str(int(opacity * 100000)))
    except Exception:  # noqa: BLE001
        pass


# Các phần tử fill có thể tồn tại trong <p:bgPr>/<a:spPr> cần dọn trước khi chèn gradFill.
_FILL_TAGS = ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill")


def _css_angle_to_drawml(angle_css_deg: float) -> int:
    """Đổi góc CSS gradient sang góc DrawingML (1/60000 độ).

    CSS: 0deg = hướng LÊN, 90deg = sang PHẢI. DrawingML ``a:lin/@ang``:
    0 = sang phải. Công thức: ang = (css_deg - 90) mod 360.
    """
    return int(round(((angle_css_deg - 90.0) % 360.0) * 60000.0))


def _apply_gradient_background(slide, angle_css: float, stops: List[Tuple[float, str]]) -> None:
    """Thay nền slide bằng gradient tuyến tính `<a:gradFill>` (thay solidFill).

    Chuỗi gradient do layout engine sinh dạng CSS; hàm này đổi sang OpenXML
    với các stop cùng vị trí %, cùng màu — đảm bảo file .pptx thật khớp
    preview web (preview dùng lại đúng chuỗi CSS đó).
    """
    fill = slide.background.fill
    fill.solid()  # đảm bảo <p:bg><p:bgPr> tồn tại, rồi ta thay phần fill bên trong

    bg = slide._element.cSld.find(qn("p:bg"))
    if bg is None:
        return
    bgPr = bg.find(qn("p:bgPr"))
    if bgPr is None:
        return

    for tag in _FILL_TAGS:
        for el in bgPr.findall(qn(tag)):
            bgPr.remove(el)

    stop_xml = "".join(
        f'<a:gs pos="{int(round(pos * 1000.0))}">'
        f'<a:srgbClr val="{hex_color.lstrip("#").upper()}"/></a:gs>'
        for pos, hex_color in stops
    )
    gradFill = parse_xml(
        f'<a:gradFill {nsdecls("a")} rotWithShape="1">'
        f"<a:gsLst>{stop_xml}</a:gsLst>"
        f'<a:lin ang="{_css_angle_to_drawml(angle_css)}" scaled="1"/>'
        f"</a:gradFill>"
    )
    bgPr.insert(0, gradFill)


def _emu_to_inches(value: Emu) -> float:
    return float(value) / EMU_PER_INCH


def _set_letter_spacing(run, spacing_pt: float) -> None:
    """Tăng khoảng cách giữa các ký tự (tracking) cho hiệu ứng 'small caps' sang trọng.

    Thuộc tính `spc` của DrawingML tính bằng 1/100 pt.
    """
    try:
        rpr = run._r.get_or_add_rPr()
        rpr.set("spc", str(int(spacing_pt * 100)))
    except Exception:  # noqa: BLE001
        pass


class PPTXMorphGeneratorService:
    """Service chuyển đổi dữ liệu PresentationResponse sang tệp .pptx hoàn chỉnh."""

    def __init__(self, data: PresentationResponse) -> None:
        self.data = data
        self.prs = Presentation()

        # Thiết lập kích thước slide theo tỉ lệ
        if self.data.aspect_ratio == "4:3":
            self.slide_width_in = 10.0
            self.slide_height_in = 7.5
        else:  # Mặc định 16:9 Widescreen
            self.slide_width_in = 13.333
            self.slide_height_in = 7.5

        self.prs.slide_width = Inches(self.slide_width_in)
        self.prs.slide_height = Inches(self.slide_height_in)
        self.blank_layout = self.prs.slide_layouts[6]

        self.palette = STYLE_PALETTES.get(
            self.data.style, STYLE_PALETTES["Futuristic"]
        )

    def _pct_to_inches(
        self, x_pct: float, y_pct: float, w_pct: float, h_pct: float
    ) -> Tuple[Emu, Emu, Emu, Emu]:
        """Quy đổi tọa độ % màn hình sang đơn vị Inches (EMU) của PowerPoint.

        Kéo hộp tràn lề về trong khung hình, ưu tiên GIỮ NGUYÊN kích thước.
        """
        w = max(0.05, min(w_pct, 100.0))
        h = max(0.05, min(h_pct, 100.0))
        x = min(max(0.0, x_pct), max(0.0, 100.0 - w))
        y = min(max(0.0, y_pct), max(0.0, 100.0 - h))

        left = Inches((x / 100.0) * self.slide_width_in)
        top = Inches((y / 100.0) * self.slide_height_in)
        width = Inches((w / 100.0) * self.slide_width_in)
        height = Inches((h / 100.0) * self.slide_height_in)
        return left, top, width, height

    def _set_slide_background(self, slide, bg_hex: Optional[str]) -> None:
        """Tô nền slide: gradient `<a:gradFill>` nếu bg là chuỗi CSS gradient,
        ngược lại tô màu đặc như cũ.

        Slide mở bài/kết bài nhận gradient theo màu chủ đề; slide thân bài giữ
        nền đặc (hoạ tiết mờ do decoration shapes đảm nhiệm).
        """
        gradient = parse_gradient_css(bg_hex)
        if gradient is not None:
            angle_css, stops = gradient
            _apply_gradient_background(slide, angle_css, stops)
            return
        color = hex_to_rgb(bg_hex, RGBColor(10, 12, 18))
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _write_speaker_notes(self, slide, slide_data: Slide) -> None:
        """Ghi lời dẫn người thuyết trình vào notes slide (khung Notes của PowerPoint)."""
        if not slide_data.speaker_notes:
            return
        try:
            slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes
        except Exception as exc:  # noqa: BLE001
            import logging
            logging.getLogger("exoticmorph.generator").warning(
                "Không ghi được speaker notes cho slide %d: %s",
                slide_data.slide_number, exc,
            )

    def _alignment(self, align: Optional[str]) -> PP_ALIGN:
        if align == "center":
            return PP_ALIGN.CENTER
        if align == "right":
            return PP_ALIGN.RIGHT
        return PP_ALIGN.LEFT

    @staticmethod
    def _body_font_size(title_size: int) -> float:
        """Cỡ chữ thân bài co giãn theo cỡ chữ tiêu đề để giữ phân cấp thị giác."""
        if title_size >= 40:
            return 11.5
        if title_size >= 24:
            return 12.5
        if title_size >= 18:
            return 11.2
        if title_size >= 15:
            return 10.4
        if title_size >= 13:
            return 9.8
        return 9.2

    # ==========================================================================
    # CÁC LOẠI PHẦN TỬ
    # ==========================================================================
    def _add_heading(self, slide, elem: SlideElement) -> None:
        """Tiêu đề chính — in đậm, cỡ lớn, độ tương phản cao."""
        left, top, width, height = self._pct_to_inches(elem.x, elem.y, elem.width, elem.height)
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = elem.content
        p.font.bold = True
        p.font.size = Pt(elem.font_size or 24)
        p.font.color.rgb = hex_to_rgb(elem.text_color, RGBColor(255, 255, 255))
        p.font.name = "Segoe UI"
        set_shape_morph_id(tx_box, elem.morph_id)

    def _add_kicker(self, slide, elem: SlideElement) -> None:
        """Nhãn mục (kicker) in hoa, có tracking — thể hiện vị trí trong mạch kể chuyện."""
        left, top, width, height = self._pct_to_inches(elem.x, elem.y, elem.width, elem.height)
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = elem.content.upper()
        run.font.bold = True
        run.font.size = Pt(elem.font_size or 12)
        run.font.color.rgb = hex_to_rgb(elem.text_color, RGBColor(200, 210, 230))
        run.font.name = "Segoe UI"
        _set_letter_spacing(run, 2.4)
        set_shape_morph_id(tx_box, elem.morph_id)

    def _add_footer(self, slide, elem: SlideElement) -> None:
        """Footer nhỏ — đồng bộ bộ nhận diện ở mọi slide."""
        left, top, width, height = self._pct_to_inches(elem.x, elem.y, elem.width, elem.height)
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = elem.content
        p.alignment = self._alignment(elem.align)
        p.font.size = Pt(elem.font_size or 9)
        p.font.color.rgb = hex_to_rgb(elem.text_color, RGBColor(150, 160, 175))
        p.font.name = "Segoe UI"
        set_shape_morph_id(tx_box, elem.morph_id)

    def _add_badge(self, slide, elem: SlideElement) -> None:
        """Pill tag / vòng tròn số thứ tự (oval khi shape_type='oval')."""
        left, top, width, height = self._pct_to_inches(elem.x, elem.y, elem.width, elem.height)
        mso = MSO_SHAPE.OVAL if elem.shape_type == "oval" else MSO_SHAPE.ROUNDED_RECTANGLE
        shape = slide.shapes.add_shape(mso, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(elem.bg_color, RGBColor(139, 92, 246))
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = elem.content
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(elem.font_size or 11)
        p.font.color.rgb = hex_to_rgb(elem.text_color, RGBColor(255, 255, 255))
        p.font.name = "Segoe UI"
        set_shape_morph_id(shape, elem.morph_id)

    def _add_accent(self, slide, elem: SlideElement) -> None:
        """Hình trang trí (đường nối lộ trình) — tô màu, không chữ."""
        left, top, width, height = self._pct_to_inches(elem.x, elem.y, elem.width, elem.height)
        mso = MSO_SHAPE.OVAL if elem.shape_type == "oval" else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(mso, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(elem.bg_color, RGBColor(70, 78, 100))
        shape.line.fill.background()
        if elem.morph_id:
            set_shape_morph_id(shape, elem.morph_id)

    def _add_decoration(self, slide, elem: SlideElement) -> None:
        """Hoạ tiết nền rất mờ (chấm sao / lưới / lattice) — opacity 3-6%.

        Vẽ TRƯỚC mọi phần tử khác (layout engine đặt decoration đầu danh sách)
        nên luôn nằm dưới card/chữ. Không có morph_id → không tham gia Morph.
        """
        left, top, width, height = self._pct_to_inches(elem.x, elem.y, elem.width, elem.height)
        mso = MSO_SHAPE.OVAL if elem.shape_type == "oval" else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(mso, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(elem.bg_color, RGBColor(120, 130, 160))
        _set_fill_opacity(shape, elem.opacity if elem.opacity is not None else 0.05)
        shape.line.fill.background()
        # Vô hiệu hóa chọn/nhấn nhầm hoạ tiết khi người dùng chỉnh sửa trong PowerPoint.
        try:
            shape._element.nvSpPr.cNvPr.set("name", f"deco_{elem.id}")
        except Exception:  # noqa: BLE001
            pass

    def _add_card(self, slide, elem: SlideElement, is_primary: bool) -> None:
        """Thẻ nội dung 2 phần: tiêu đề (thông điệp cụ thể) + đoạn diễn giải chi tiết.

        Thiết kế + biến thể thị giác (phá khuôn card đồng nhất):
        - ``accent_top`` (mặc định): nền đặc + thanh accent trên đỉnh.
        - ``accent_left``: nền đặc + VIỀN TRÁI DÀY (~4.5px) thay thanh trên —
          nhịp điệu riêng cho card phụ / list-rows.
        - ``outline``: KHÔNG nền (trong suốt) + viền mỏng màu accent + icon cạnh
          tiêu đề → phân cấp chính/phụ rõ ràng thay vì đồng trọng lượng.
        - ``cta``: thanh Call-To-Action tô đặc toàn thân.
        - Card icon (``elem.icon``): icon motif chủ đề đặt cạnh tiêu đề; card số
          lớn có icon xu hướng/góc kèm nhãn ở góc trên phải.
        - BIG_STAT_CALLOUT dùng text box riêng, ``word_wrap=False`` cho số chính
          để chuỗi như ``99.86%`` luôn nằm trên một dòng.
        """
        left, top, width, height = self._pct_to_inches(elem.x, elem.y, elem.width, elem.height)
        accent_rgb = hex_to_rgb(elem.accent_color, RGBColor(16, 185, 129))
        title_size = int(elem.font_size or 15)
        is_stat_card = title_size >= 40
        is_cta = (elem.variant == "cta") or (elem.align == "center" and not elem.sub_text)
        variant = elem.variant or "accent_top"

        # --- Thân thẻ (vẽ trước, nằm dưới thanh accent / viền trái) ---
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        if variant == "outline":
            # Card "không nền": trong suốt + viền mỏng màu accent.
            shape.fill.background()
            shape.line.color.rgb = accent_rgb
            shape.line.width = Pt(OUTLINE_BORDER_PT)
            _set_line_opacity(shape, OUTLINE_BORDER_OPACITY)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(elem.bg_color, RGBColor(20, 24, 39))
            if elem.border_color:
                shape.line.color.rgb = hex_to_rgb(elem.border_color, RGBColor(40, 46, 66))
                shape.line.width = Pt(CARD_BORDER_PT)
                _set_line_opacity(shape, CARD_BORDER_OPACITY)
            else:
                shape.line.fill.background()

        # Thanh accent: CHỈ variant accent_top (và card thường) mới có thanh trên đỉnh.
        if variant == "accent_top" and not is_cta and elem.type != "shape":
            bar_h = 0.075 if is_primary else 0.045
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left + Inches(0.08), top + Inches(0.045), width - Inches(0.16), Inches(bar_h),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_rgb
            bar.line.fill.background()
            if elem.morph_id:
                set_shape_morph_id(bar, f"{elem.morph_id}_accent")

        # Biến thể accent_left: viền trái dày thay cho thanh trên đỉnh.
        if variant == "accent_left" and not is_cta and elem.type != "shape":
            bar_v_inset = Inches(0.10)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left + Inches(0.02),
                top + bar_v_inset,
                Inches(CARD_LEFT_BAR_IN),
                height - 2 * bar_v_inset,
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_rgb
            bar.line.fill.background()
            if elem.morph_id:
                set_shape_morph_id(bar, f"{elem.morph_id}_accent")

        # Icon motif cho card SỐ KHỔNG LỒ: icon xu hướng/đơn vị ở góc trên phải.
        if is_stat_card and elem.icon:
            icon_box = slide.shapes.add_textbox(
                left + width - Inches(0.62), top + Inches(0.16),
                Inches(0.44), Inches(0.36),
            )
            itf = icon_box.text_frame
            itf.word_wrap = False
            itf.margin_left = itf.margin_top = itf.margin_right = itf.margin_bottom = 0
            ip = itf.paragraphs[0]
            ip.alignment = PP_ALIGN.RIGHT
            ir = ip.add_run()
            ir.text = elem.icon
            ir.font.size = Pt(16)
            ir.font.bold = True
            ir.font.color.rgb = accent_rgb
            # Không gắn font name: để PowerPoint tự chọn Segoe UI Emoji cho emoji.
            if elem.morph_id:
                set_shape_morph_id(icon_box, f"{elem.morph_id}_icon")

        if is_stat_card:
            # Shape chỉ giữ nền/viền; text được đặt bằng 2 textbox để số chính
            # không wrap còn mô tả vẫn wrap bình thường.
            shape.text_frame.clear()
            w_in = _emu_to_inches(width)
            h_in = _emu_to_inches(height)
            pad = CARD_TEXT_PAD_IN
            inner_w = max(0.2, w_in - 2 * pad)
            title_h = min(1.14, max(0.72, h_in * 0.34))
            body_gap = 0.16
            body_h = max(0.30, h_in - title_h - body_gap - 0.74)
            title_y = 0.54 if h_in >= 2.2 else CARD_TEXT_PAD_TOP_IN
            body_y = min(h_in - body_h - CARD_TEXT_PAD_BOTTOM_IN, title_y + title_h + body_gap)

            title_box = slide.shapes.add_textbox(
                left + Inches(pad), top + Inches(title_y), Inches(inner_w), Inches(title_h)
            )
            title_tf = title_box.text_frame
            title_tf.word_wrap = False
            title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            title_tf.margin_left = title_tf.margin_right = 0
            title_tf.margin_top = title_tf.margin_bottom = 0
            p_title = title_tf.paragraphs[0]
            p_title.text = elem.content
            p_title.alignment = PP_ALIGN.CENTER
            p_title.font.bold = True
            p_title.font.size = Pt(title_size)
            p_title.font.color.rgb = hex_to_rgb(elem.text_color, RGBColor(255, 255, 255))
            p_title.font.name = "Segoe UI"
            set_shape_morph_id(title_box, f"{elem.morph_id}_value")

            if elem.sub_text:
                body_box = slide.shapes.add_textbox(
                    left + Inches(pad), top + Inches(body_y), Inches(inner_w), Inches(body_h)
                )
                body_tf = body_box.text_frame
                body_tf.word_wrap = True
                body_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                body_tf.margin_left = body_tf.margin_right = 0
                body_tf.margin_top = body_tf.margin_bottom = 0
                p_body = body_tf.paragraphs[0]
                p_body.text = elem.sub_text.strip()
                p_body.alignment = PP_ALIGN.CENTER
                p_body.font.size = Pt(self._body_font_size(title_size))
                p_body.font.color.rgb = hex_to_rgb(self.palette.get("sub_text"), RGBColor(165, 175, 195))
                p_body.font.name = "Segoe UI"
                p_body.line_spacing = 1.12
                set_shape_morph_id(body_box, f"{elem.morph_id}_body")

            set_shape_morph_id(shape, elem.morph_id)
            return

        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(CARD_TEXT_PAD_IN + (CARD_LEFT_BAR_IN if variant == "accent_left" else 0.0))
        tf.margin_right = Inches(CARD_TEXT_PAD_IN)
        tf.margin_top = Inches(CARD_TEXT_PAD_TOP_IN if not is_cta else 0.08)
        tf.margin_bottom = Inches(CARD_TEXT_PAD_BOTTOM_IN if not is_cta else 0.08)
        try:
            # Canh giữa dọc cho MỌI card (không chỉ CTA): khi nội dung ngắn hơn
            # khung (khung được cấp theo hằng số cố định, không co theo chữ),
            # văn bản sẽ nằm giữa thay vì dồn lên đỉnh và để trống hẳn phía dưới —
            # giảm cảm giác "khung to, chữ bé, lệch hẳn lên trên".
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:  # noqa: BLE001
            pass

        # 1) Tiêu đề — thông điệp cụ thể, in đậm, tương phản cao.
        #    Card có icon motif: icon là run riêng (không gắn font name để emoji
        #    fallback sang Segoe UI Emoji trong PowerPoint).
        p_title = tf.paragraphs[0]
        p_title.alignment = self._alignment(elem.align)
        if elem.icon:
            icon_run = p_title.add_run()
            icon_run.text = f"{elem.icon}  "
            icon_run.font.bold = True
            icon_run.font.size = Pt(title_size)
            icon_run.font.color.rgb = accent_rgb
        title_run = p_title.add_run()
        title_run.text = elem.content
        title_run.font.bold = True
        title_run.font.size = Pt(title_size)
        title_run.font.color.rgb = hex_to_rgb(elem.text_color, RGBColor(255, 255, 255))
        title_run.font.name = "Segoe UI"

        # 2) Đoạn diễn giải chi tiết — màu phụ, khoảng cách dòng thoáng.
        if elem.sub_text:
            body = elem.sub_text.strip()
            p_body = tf.add_paragraph()
            p_body.space_before = Pt(5)
            p_body.text = body
            p_body.alignment = self._alignment(elem.align)
            p_body.font.size = Pt(self._body_font_size(title_size))
            p_body.font.color.rgb = hex_to_rgb(self.palette.get("sub_text"), RGBColor(165, 175, 195))
            p_body.font.name = "Segoe UI"
            p_body.line_spacing = 1.12

        set_shape_morph_id(shape, elem.morph_id)

    def _add_element_to_slide(self, slide, elem: SlideElement, is_primary: bool = False) -> None:
        """Phân luồng theo loại phần tử."""
        if elem.type == "heading":
            self._add_heading(slide, elem)
        elif elem.type == "kicker":
            self._add_kicker(slide, elem)
        elif elem.type == "footer":
            self._add_footer(slide, elem)
        elif elem.type == "badge":
            self._add_badge(slide, elem)
        elif elem.type == "accent":
            self._add_accent(slide, elem)
        elif elem.type == "decoration":
            self._add_decoration(slide, elem)
        else:
            # card / metric / text / shape mặc định dùng kiểu thẻ nội dung.
            self._add_card(slide, elem, is_primary)

    def build_presentation(self) -> io.BytesIO:
        """Dựng toàn bộ slide, chèn Morph XML transitions và xuất ra BytesIO."""
        if not self.data.slides:
            raise ValueError("Dữ liệu bài trình chiếu trống: không có slide nào để dựng file PPTX.")

        for slide_idx, slide_data in enumerate(self.data.slides):
            slide = self.prs.slides.add_slide(self.blank_layout)
            self._set_slide_background(slide, slide_data.bg_color)
            self._write_speaker_notes(slide, slide_data)

            # Xác định thẻ chính (primary) = thẻ đầu tiên trong danh sách (nếu có >= 2 thẻ).
            card_elements = [e for e in slide_data.elements if e.type in ("card", "text", "shape", "metric")]
            primary_id = card_elements[0].id if len(card_elements) >= 2 else None

            for elem in slide_data.elements:
                self._add_element_to_slide(slide, elem, is_primary=(elem.id == primary_id))

            # Kích hoạt hiệu ứng Morph cho tất cả slide từ slide 2 trở đi.
            if slide_idx >= 1:
                enable_morph_transition(
                    slide=slide,
                    duration_ms=MORPH_DURATION_MS,
                    morph_option="byObject",
                    advance_on_click=True,
                )

        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer


def build_pptx_from_schema(data: PresentationResponse) -> io.BytesIO:
    """Hàm helper tiện ích dựng file PPTX từ PresentationResponse."""
    service = PPTXMorphGeneratorService(data)
    return service.build_presentation()
